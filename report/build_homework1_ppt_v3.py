"""Generate a presentation for homework 1.

The deck is built from the existing code outputs and recomputed metrics so the
slides stay consistent with the project implementation.
"""

from __future__ import annotations

import sys
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from homework1.backtest import run_backtest
from homework1.config import (
    BUY_THRESHOLD,
    COMMISSION,
    INITIAL_CAPITAL,
    SELL_THRESHOLD,
    WINDOW,
)
from homework1.features import build_features
from homework1.models import evaluate_model, get_models


OUT_DIR = ROOT / "outputs" / "homework1"
FIG_DIR = OUT_DIR / "figures"
PPTX_PATH = OUT_DIR / "homework1_presentation_v3.pptx"
ASSET_DIR = OUT_DIR / "ppt_assets"

W, H = 13.333, 7.5

COLORS = {
    "ink": "202124",
    "muted": "667085",
    "soft": "F5F7FA",
    "paper": "FFFFFF",
    "green": "1B7F5C",
    "mint": "DDF3EA",
    "gold": "DFAF2B",
    "red": "B5403C",
    "navy": "22304D",
    "line": "D8DEE8",
    "blue": "2F6FED",
    "orange": "D97904",
    "purple": "6E56CF",
}

MODEL_COLORS = {
    "LinearRegression": COLORS["blue"],
    "RandomForest": COLORS["orange"],
    "LightGBM": COLORS["green"],
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    color: str = COLORS["ink"],
    bold: bool = False,
    font: str = "Microsoft YaHei",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_lines(slide, lines, x, y, w, h, *, size=16, color=COLORS["ink"], line_spacing=1.1):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(5 * line_spacing)
    return box


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = COLORS["paper"],
    line: str | None = None,
    radius: bool = False,
    transparency: int = 0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_title(slide, title: str, subtitle: str | None = None, *, inverse: bool = False):
    title_color = COLORS["paper"] if inverse else COLORS["ink"]
    sub_color = "D7E2D8" if inverse else COLORS["muted"]
    add_text(slide, title, 0.65, 0.36, 8.6, 0.55, size=27, bold=True, color=title_color, margin=0)
    if subtitle:
        add_text(slide, subtitle, 0.67, 0.95, 9.2, 0.32, size=10, color=sub_color, margin=0)


def add_footer(slide, page: int, inverse: bool = False):
    c = "D7E2D8" if inverse else "98A2B3"
    add_text(slide, "金融与经济数据挖掘 · 作业1", 0.65, 7.16, 4.2, 0.22, size=8, color=c, margin=0)
    add_text(slide, f"{page:02d}", 12.35, 7.16, 0.35, 0.22, size=8, color=c, align=PP_ALIGN.RIGHT, margin=0)


def add_kpi(slide, label: str, value: str, x: float, y: float, w: float, h: float, *, color: str, value_size: int = 27):
    add_rect(slide, x, y, w, h, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, value, x + 0.18, y + 0.18, w - 0.36, 0.45, size=value_size, bold=True, color=color, margin=0)
    add_text(slide, label, x + 0.2, y + 0.78, w - 0.38, 0.28, size=10, color=COLORS["muted"], margin=0)


def add_metric_card(slide, model: str, ret: str, dd: str, wr: str, trades: str, x: float, y: float, w: float, h: float):
    color = MODEL_COLORS[model]
    add_rect(slide, x, y, w, h, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_rect(slide, x, y, w, 0.08, fill=color, radius=False)
    add_text(slide, model, x + 0.22, y + 0.22, w - 0.44, 0.28, size=15, bold=True, color=COLORS["ink"], margin=0)
    add_text(slide, ret, x + 0.22, y + 0.60, w - 0.44, 0.55, size=26, bold=True, color=color, margin=0)
    add_text(slide, "累计收益", x + 0.24, y + 1.10, w - 0.44, 0.22, size=9, color=COLORS["muted"], margin=0)
    add_text(slide, f"最大回撤  {dd}", x + 0.24, y + 1.48, w - 0.44, 0.22, size=10, color=COLORS["ink"], margin=0)
    add_text(slide, f"胜率  {wr}", x + 0.24, y + 1.83, w - 0.44, 0.22, size=10, color=COLORS["ink"], margin=0)
    add_text(slide, f"交易次数  {trades}", x + 0.24, y + 2.18, w - 0.44, 0.22, size=10, color=COLORS["ink"], margin=0)


def add_image_fit(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    box_ratio = w / h
    if ratio > box_ratio:
        new_w = w
        new_h = w / ratio
        dx = 0
        dy = (h - new_h) / 2
    else:
        new_h = h
        new_w = h * ratio
        dx = (w - new_w) / 2
        dy = 0
    return slide.shapes.add_picture(str(path), Inches(x + dx), Inches(y + dy), Inches(new_w), Inches(new_h))


def add_code_block(
    slide,
    title: str,
    code: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 9,
    title_color: str = COLORS["gold"],
):
    """Draw a dark code panel with a small file/title label."""
    add_rect(slide, x, y, w, h, fill=COLORS["navy"], line="32415F", radius=True)
    add_rect(slide, x, y, w, 0.50, fill="18243A", radius=False)
    add_text(slide, title, x + 0.22, y + 0.14, w - 0.44, 0.22, size=10, bold=True, color=title_color, font="Consolas", margin=0)
    add_text(
        slide,
        code.strip(),
        x + 0.28,
        y + 0.65,
        w - 0.56,
        h - 0.85,
        size=size,
        color="E7ECF8",
        font="Consolas",
        margin=0,
    )


def set_background(slide, color: str):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(color)


def compute_results():
    df = pd.read_csv(ROOT / "data" / "homework1" / "sh_600519.csv", parse_dates=["date"]).set_index("date")
    df_feat, feature_cols = build_features(df, window=WINDOW)

    df_2024 = df_feat[df_feat.index < "2025-01-01"]
    df_2025 = df_feat[df_feat.index >= "2025-01-01"]
    x_2024 = df_2024[feature_cols]
    y_2024 = df_2024["label"]
    split_idx = int(len(df_2024) * 0.8)
    x_train = x_2024.iloc[:split_idx]
    x_test = x_2024.iloc[split_idx:]
    y_train = y_2024.iloc[:split_idx]
    y_test = y_2024.iloc[split_idx:]

    models = get_models()
    metrics = {}
    backtests = {}
    preds_2025 = {}

    for name, model in models.items():
        model.fit(x_train, y_train)
        metrics[name] = evaluate_model(y_test.values, model.predict(x_test))

        final_model = get_models()[name]
        final_model.fit(x_2024, y_2024)
        pred = final_model.predict(df_2025[feature_cols])
        preds_2025[name] = pred
        backtests[name] = run_backtest(
            df_2025,
            pred,
            initial_capital=INITIAL_CAPITAL,
            commission=COMMISSION,
            buy_threshold=BUY_THRESHOLD,
            sell_threshold=SELL_THRESHOLD,
        )

    return {
        "df": df,
        "df_feat": df_feat,
        "df_2024": df_2024,
        "df_2025": df_2025,
        "split_idx": split_idx,
        "x_train": x_train,
        "x_test": x_test,
        "metrics": metrics,
        "backtests": backtests,
        "preds_2025": preds_2025,
    }


def pct(v: float, digits: int = 2) -> str:
    return f"{v * 100:.{digits}f}%"


def money(v: float) -> str:
    return f"{v / 10000:.1f}万"


def create_support_charts(result: dict):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    df = result["df"]
    df_2025 = result["df_2025"]
    metrics = result["metrics"]
    backtests = result["backtests"]

    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False

    # Close and volume overview.
    fig, axes = plt.subplots(2, 1, figsize=(11, 5.4), sharex=True, gridspec_kw={"height_ratios": [2, 1]})
    axes[0].plot(df.index, df["close"], color=f"#{COLORS['green']}", linewidth=1.8)
    axes[0].axvspan(pd.Timestamp("2025-01-01"), df.index.max(), color=f"#{COLORS['mint']}", alpha=0.55)
    axes[0].set_ylabel("前复权收盘价")
    axes[0].grid(True, alpha=0.18)
    axes[1].bar(df.index, df["volume"] / 10000, color=f"#{COLORS['gold']}", width=1.0, alpha=0.75)
    axes[1].axvspan(pd.Timestamp("2025-01-01"), df.index.max(), color=f"#{COLORS['mint']}", alpha=0.55)
    axes[1].set_ylabel("成交量(万股)")
    axes[1].grid(True, alpha=0.16)
    fig.suptitle("贵州茅台 2024-2025 收盘价与成交量", fontsize=15, fontweight="bold")
    plt.tight_layout()
    overview_path = ASSET_DIR / "data_overview.png"
    fig.savefig(overview_path, dpi=180, bbox_inches="tight")
    plt.close(fig)

    # ML metrics comparison.
    models = list(metrics.keys())
    rmse = [metrics[m]["RMSE"] for m in models]
    mae = [metrics[m]["MAE"] for m in models]
    y = np.arange(len(models))
    fig, ax = plt.subplots(figsize=(8.2, 4.1))
    ax.barh(y + 0.18, rmse, height=0.32, label="RMSE", color=f"#{COLORS['navy']}")
    ax.barh(y - 0.18, mae, height=0.32, label="MAE", color=f"#{COLORS['green']}")
    ax.set_yticks(y)
    ax.set_yticklabels(models)
    ax.invert_yaxis()
    ax.set_xlabel("误差")
    ax.set_title("2024 测试集机器学习指标", fontsize=14, fontweight="bold")
    ax.grid(axis="x", alpha=0.18)
    ax.legend(frameon=False, loc="upper right")
    for i, v in enumerate(rmse):
        ax.text(v + 0.00025, i + 0.18, f"{v:.4f}", va="center", fontsize=9)
    for i, v in enumerate(mae):
        ax.text(v + 0.00025, i - 0.18, f"{v:.4f}", va="center", fontsize=9)
    plt.tight_layout()
    metric_path = ASSET_DIR / "metrics_bar.png"
    fig.savefig(metric_path, dpi=180, bbox_inches="tight")
    plt.close(fig)

    # Return and drawdown comparison.
    returns = [backtests[m]["total_return"] for m in models]
    drawdowns = [-backtests[m]["max_drawdown"] for m in models]
    fig, ax = plt.subplots(figsize=(8.2, 4.3))
    width = 0.34
    x = np.arange(len(models))
    bars1 = ax.bar(x - width / 2, returns, width=width, label="累计收益", color=[f"#{MODEL_COLORS[m]}" for m in models])
    bars2 = ax.bar(x + width / 2, drawdowns, width=width, label="最大回撤(向下)", color="#AAB2C0")
    ax.axhline(0, color="#394150", linewidth=0.9)
    ax.set_xticks(x)
    ax.set_xticklabels(models)
    ax.set_ylabel("比例")
    ax.set_title("2025 回测收益与风险", fontsize=14, fontweight="bold")
    ax.grid(axis="y", alpha=0.16)
    ax.legend(frameon=False, loc="upper left")
    for bar in bars1:
        v = bar.get_height()
        ax.text(bar.get_x() + bar.get_width() / 2, v + (0.008 if v >= 0 else -0.015), pct(v), ha="center", va="bottom" if v >= 0 else "top", fontsize=8)
    offsets = [-0.022, -0.028, -0.015]
    for idx, bar in enumerate(bars2):
        v = bar.get_height()
        ax.text(bar.get_x() + bar.get_width() / 2, v + offsets[idx], pct(abs(v)), ha="center", va="top", fontsize=7.5, color="#394150")
    plt.tight_layout()
    ret_path = ASSET_DIR / "return_drawdown.png"
    fig.savefig(ret_path, dpi=180, bbox_inches="tight")
    plt.close(fig)

    # Buy-and-hold comparison for 2025.
    buy_hold = df[df.index.year == 2025]["close"].iloc[-1] / df[df.index.year == 2025]["close"].iloc[0] - 1
    labels = ["买入持有", "LightGBM策略"]
    vals = [buy_hold, backtests["LightGBM"]["total_return"]]
    fig, ax = plt.subplots(figsize=(5.4, 3.7))
    bars = ax.bar(labels, vals, color=[f"#{COLORS['red']}", f"#{COLORS['green']}"], width=0.54)
    ax.axhline(0, color="#394150", linewidth=0.9)
    ax.set_title("2025 策略价值对照", fontsize=14, fontweight="bold")
    ax.set_ylabel("累计收益")
    ax.grid(axis="y", alpha=0.16)
    for bar, v in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width() / 2, v + (0.008 if v >= 0 else -0.012), pct(v), ha="center", va="bottom" if v >= 0 else "top", fontsize=11, fontweight="bold")
    plt.tight_layout()
    bh_path = ASSET_DIR / "buy_hold_compare.png"
    fig.savefig(bh_path, dpi=180, bbox_inches="tight")
    plt.close(fig)

    return {
        "overview": overview_path,
        "metrics": metric_path,
        "return_drawdown": ret_path,
        "buy_hold": bh_path,
    }


def add_table(slide, rows, x, y, w, h, col_widths, *, header_fill=COLORS["green"]):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for i, cw in enumerate(col_widths):
        table.columns[i].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            fill = header_fill if r == 0 else (COLORS["soft"] if r % 2 == 0 else COLORS["paper"])
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(fill)
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c != 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(10 if r else 10)
                    run.font.bold = r == 0 or (value == "LightGBM")
                    run.font.color.rgb = rgb(COLORS["paper"] if r == 0 else COLORS["ink"])
    return table_shape


def build_deck(result: dict, charts: dict):
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    df = result["df"]
    df_2024 = result["df_2024"]
    df_2025 = result["df_2025"]
    metrics = result["metrics"]
    backtests = result["backtests"]
    buy_hold_2025 = df[df.index.year == 2025]["close"].iloc[-1] / df[df.index.year == 2025]["close"].iloc[0] - 1

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_rect(slide, 0, 0, 0.25, 7.5, fill=COLORS["green"])
    add_text(slide, "金融与经济数据挖掘 · 作业 1", 0.72, 0.55, 5.5, 0.35, size=13, color=COLORS["muted"], margin=0)
    add_text(slide, "机器学习技术分析实战", 0.72, 1.2, 8.5, 0.75, size=36, bold=True, color=COLORS["ink"], margin=0)
    add_text(slide, "贵州茅台 sh.600519 · 2024-2025 日线 · 三模型对比与回测", 0.74, 2.1, 8.0, 0.35, size=13, color=COLORS["muted"], margin=0)
    add_rect(slide, 0.72, 2.8, 11.85, 0.02, fill=COLORS["line"])
    add_kpi(slide, "研究标的", "贵州茅台", 0.75, 3.25, 2.55, 1.15, color=COLORS["green"])
    add_kpi(slide, "特征窗口", "10日", 3.55, 3.25, 1.65, 1.15, color=COLORS["navy"])
    add_kpi(slide, "初始资金", "100万", 5.45, 3.25, 1.75, 1.15, color=COLORS["gold"])
    add_rect(slide, 7.65, 3.25, 4.92, 3.25, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_rect(slide, 7.65, 3.25, 4.92, 0.52, fill=COLORS["green"])
    add_text(slide, "核心结果", 7.92, 3.38, 2.0, 0.28, size=13, bold=True, color=COLORS["paper"], margin=0)
    add_text(slide, "+10.40%", 7.92, 4.0, 3.0, 0.55, size=32, bold=True, color=COLORS["green"], margin=0)
    add_text(slide, "LightGBM 2025 年回测累计收益", 7.95, 4.58, 3.8, 0.22, size=10, color=COLORS["muted"], margin=0)
    add_text(slide, "5.94%", 7.92, 5.15, 2.0, 0.42, size=24, bold=True, color=COLORS["navy"], margin=0)
    add_text(slide, "最大回撤，三组模型中最低", 7.95, 5.6, 3.0, 0.2, size=9, color=COLORS["muted"], margin=0)
    add_text(slide, "54.96%", 10.65, 5.15, 1.6, 0.42, size=24, bold=True, color=COLORS["gold"], margin=0)
    add_text(slide, "方向胜率", 10.68, 5.6, 1.5, 0.2, size=9, color=COLORS["muted"], margin=0)
    add_footer(slide, 1)

    # 2. Assignment requirements
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "作业要求拆解", "从原始行情数据到策略回测结论，形成一条完整分析链路")
    x0, y0 = 0.82, 1.42
    steps = [
        ("01", "数据获取", "baostock 日线数据\n前复权收盘价 + 成交量"),
        ("02", "特征构建", "价格收益率 + 成交量收益率\n过去 10 日滚动窗口"),
        ("03", "模型训练", "Linear / RandomForest / LightGBM\n2024 年时序切分评估"),
        ("04", "实战回测", "全部 2024 年训练最终模型\n2025 年执行买卖规则"),
        ("05", "结果输出", "RMSE、MAE、累计收益\n最大回撤、胜率、交易次数"),
        ("06", "课堂展示", "步骤、代码、图表、结论\n对应作业提交要求"),
    ]
    for i, (num, name, desc) in enumerate(steps):
        row, col = divmod(i, 3)
        x = x0 + col * 4.08
        y = y0 + row * 1.95
        add_rect(slide, x, y, 3.62, 1.48, fill=COLORS["paper"], line=COLORS["line"], radius=True)
        add_rect(slide, x + 0.22, y + 0.25, 0.54, 0.42, fill=COLORS["green"], radius=True)
        add_text(slide, num, x + 0.31, y + 0.34, 0.36, 0.12, size=9, bold=True, color=COLORS["paper"], align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, name, x + 0.95, y + 0.25, 2.1, 0.3, size=18, bold=True, color=COLORS["ink"], margin=0)
        add_rich_lines(slide, desc.split("\n"), x + 0.95, y + 0.72, 2.3, 0.42, size=10, color=COLORS["muted"], line_spacing=0.75)
    add_rect(slide, 0.9, 5.68, 11.55, 0.62, fill=COLORS["mint"], line=COLORS["line"], radius=True)
    add_text(slide, "展示逻辑：先说明为什么这样建模，再用代码片段证明实现，最后用图表和指标解释投资价值。", 1.18, 5.9, 11.0, 0.16, size=12, bold=True, color=COLORS["green"], align=PP_ALIGN.CENTER, margin=0)
    add_footer(slide, 2)

    # 3. Data overview
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "数据说明与样本质量", "研究标的：贵州茅台 sh.600519，A 股白酒行业龙头")
    add_image_fit(slide, charts["overview"], 0.73, 1.42, 7.25, 4.1)
    # KPI cards in 2x2 grid for better spacing
    add_kpi(slide, "原始交易日", f"{len(df)}天", 8.28, 1.42, 2.35, 1.03, color=COLORS["green"])
    add_kpi(slide, "2024 / 2025", f"{(df.index.year == 2024).sum()} / {(df.index.year == 2025).sum()}", 10.88, 1.42, 2.15, 1.03, color=COLORS["navy"])
    add_kpi(slide, "价格区间", f"{df.close.min():.0f} - {df.close.max():.0f}", 8.28, 2.65, 2.35, 1.03, color=COLORS["gold"], value_size=20)
    add_kpi(slide, "成交量区间", "103-1947万股", 10.88, 2.65, 2.15, 1.03, color=COLORS["green"], value_size=20)
    add_rect(slide, 8.28, 3.95, 4.75, 1.55, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "数据质量检查", 8.55, 4.15, 2.2, 0.27, size=15, bold=True, color=COLORS["ink"], margin=0)
    add_rich_lines(slide, ["无缺失值；日期索引升序", "close 转为浮点数，volume 转为数值", "前复权口径减少除权除息干扰"], 8.55, 4.55, 4.15, 0.7, size=11, color=COLORS["muted"], line_spacing=0.9)
    add_footer(slide, 3)

    # 4. Feature engineering
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "特征工程：用变动预测变动", "不直接使用价格水平，而使用收益率序列表达短期变化")
    add_rect(slide, 0.74, 1.35, 5.75, 4.85, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "计算公式", 1.05, 1.68, 1.8, 0.3, size=18, bold=True, color=COLORS["green"], margin=0)
    add_text(slide, "price_returnₜ = (closeₜ - closeₜ₋₁) / closeₜ₋₁", 1.05, 2.22, 4.85, 0.32, size=15, bold=True, color=COLORS["ink"], margin=0)
    add_text(slide, "vol_returnₜ = (volumeₜ - volumeₜ₋₁) / volumeₜ₋₁", 1.05, 2.88, 4.9, 0.32, size=15, bold=True, color=COLORS["ink"], margin=0)
    add_text(slide, "labelₜ = price_returnₜ₊₁", 1.05, 3.54, 3.4, 0.32, size=15, bold=True, color=COLORS["ink"], margin=0)
    add_rich_lines(slide, ["标签含义：预测下一交易日价格收益率", "滚动窗口：过去 10 日价格收益率 + 过去 10 日成交量收益率", "特征总数：20 个"], 1.05, 4.35, 4.75, 0.78, size=12, color=COLORS["muted"], line_spacing=0.95)
    add_rect(slide, 7.0, 1.35, 5.58, 4.85, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "样本构造结果", 7.32, 1.68, 2.4, 0.3, size=18, bold=True, color=COLORS["green"], margin=0)
    # KPI cards in 2x2 grid for better spacing
    add_kpi(slide, "特征后样本", f"{len(result['df_feat'])}条", 7.35, 2.28, 2.45, 1.0, color=COLORS["green"])
    add_kpi(slide, "2024样本", f"{len(df_2024)}条", 10.05, 2.28, 2.25, 1.0, color=COLORS["navy"])
    add_kpi(slide, "2025样本", f"{len(df_2025)}条", 7.35, 3.45, 2.45, 1.0, color=COLORS["gold"])
    add_text(slide, "去除缺失值来源", 10.05, 3.65, 2.2, 0.26, size=14, bold=True, color=COLORS["ink"], margin=0)
    add_rich_lines(slide, ["前 10 日无法形成完整滞后窗口", "最后 1 日没有次日收益率标签", "因此 485 个原始交易日生成 473 条样本"], 10.05, 4.05, 2.25, 0.76, size=10, color=COLORS["muted"], line_spacing=0.9)
    add_footer(slide, 4)

    # 5. Core code: feature construction
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "核心代码：特征构建", "对应 homework1/features.py，完成收益率、标签与 10 日滚动窗口")
    feature_code = """
df["price_return"] = df["close"].pct_change()
df["vol_return"] = df["volume"].pct_change()
df["label"] = df["price_return"].shift(-1)

for i in range(1, window + 1):
    df[f"price_return_{i}"] = df["price_return"].shift(i)
    df[f"vol_return_{i}"] = df["vol_return"].shift(i)

df = df.dropna()
feature_cols = (
    [f"price_return_{i}" for i in range(1, window + 1)]
    + [f"vol_return_{i}" for i in range(1, window + 1)]
)
"""
    add_code_block(slide, "homework1/features.py", feature_code, 0.82, 1.38, 7.45, 4.85, size=13)
    add_rect(slide, 8.58, 1.38, 4.0, 4.85, fill=COLORS["soft"], line=COLORS["line"], radius=True)
    add_text(slide, "代码讲解", 8.88, 1.72, 1.6, 0.28, size=18, bold=True, color=COLORS["green"], margin=0)
    add_rich_lines(
        slide,
        [
            "pct_change() 把价格和成交量",
            "  转换为日变化率，降低量纲影响。",
            "shift(-1) 构造次日价格收益率，",
            "  作为监督学习标签。",
            "循环生成 10 日滞后序列，",
            "  最终得到 20 个输入特征。",
            "dropna() 清理窗口前端和",
            "  最后一天标签缺失的样本。",
        ],
        8.88,
        2.24,
        3.45,
        2.5,
        size=11,
        color=COLORS["muted"],
        line_spacing=0.85,
    )
    add_kpi(slide, "close / volume", "2列", 8.88, 4.72, 1.8, 1.05, color=COLORS["green"])
    add_kpi(slide, "输出特征", "20个", 10.92, 4.72, 1.3, 1.05, color=COLORS["gold"])
    add_footer(slide, 5)

    # 6. Time split and models
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "建模流程与时间切分", "2024 年评估模型，全部 2024 年训练最终模型，2025 年独立回测")
    # time split cards
    y = 1.85
    split_cards = [
        (0.95, 3.65, COLORS["mint"], COLORS["green"], "2024 前80%", "训练集 184 条", f"{result['x_train'].index.min().date()} - {result['x_train'].index.max().date()}"),
        (5.05, 2.75, "F7E6B5", COLORS["orange"], "2024 后20%", "测试集 47 条", f"{result['x_test'].index.min().date()} - {result['x_test'].index.max().date()}"),
        (8.25, 3.95, "E7ECF8", COLORS["navy"], "2025", "回测集 242 条", f"{df_2025.index.min().date()} - {df_2025.index.max().date()}"),
    ]
    for x, w, fill, color, period, role, dates in split_cards:
        add_rect(slide, x, y, w, 1.2, fill=fill, line=COLORS["line"], radius=True)
        add_text(slide, period, x + 0.26, y + 0.22, w - 0.5, 0.22, size=12, bold=True, color=color, margin=0)
        add_text(slide, role, x + 0.26, y + 0.55, w - 0.5, 0.22, size=16, bold=True, color=COLORS["ink"], margin=0)
        add_text(slide, dates, x + 0.26, y + 0.91, w - 0.5, 0.16, size=9, color=COLORS["muted"], margin=0)
    for x in [4.68, 7.88]:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y + 0.35), Inches(0.45), Inches(0.45))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = rgb(COLORS["gold"])
        arrow.line.fill.background()
    models = [
        ("LinearRegression", "线性基线", "简单、可解释、训练快"),
        ("RandomForest", "集成基线", "捕捉非线性，抗噪能力较强"),
        ("LightGBM", "主模型", "梯度提升树，回测表现最佳"),
    ]
    for i, (m, tag, desc) in enumerate(models):
        x = 0.95 + i * 4.13
        add_rect(slide, x, 4.05, 3.4, 1.32, fill=COLORS["soft"], line=COLORS["line"], radius=True)
        add_rect(slide, x + 0.18, 4.25, 0.28, 0.88, fill=MODEL_COLORS[m], radius=True)
        add_text(slide, m, x + 0.62, 4.23, 2.35, 0.25, size=15, bold=True, color=COLORS["ink"], margin=0)
        add_text(slide, tag, x + 0.62, 4.62, 1.3, 0.22, size=10, bold=True, color=MODEL_COLORS[m], margin=0)
        add_text(slide, desc, x + 0.62, 4.93, 2.45, 0.22, size=10, color=COLORS["muted"], margin=0)
    add_footer(slide, 6)

    # 7. Core code: model training and backtest signal
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "核心代码：训练评估与交易信号", "展示时序切分、模型训练和回测规则的关键实现")
    train_code = """
mask_2024 = df_feat.index < "2025-01-01"
df_2024 = df_feat[mask_2024]
split_idx = int(len(df_2024) * 0.8)

X_train = X_2024.iloc[:split_idx]
X_test = X_2024.iloc[split_idx:]
y_train = y_2024.iloc[:split_idx]
y_test = y_2024.iloc[split_idx:]

for name, model in get_models().items():
    model.fit(X_train, y_train)
    y_pred = model.predict(X_test)
"""
    signal_code = """
if pred > buy_threshold and not position:
    shares = capital * (1 - commission) / price
    capital = 0
    position = True
elif pred < -sell_threshold and position:
    capital = shares * price * (1 - commission)
    shares = 0
    position = False

daily_capital.append(capital + shares * price)
"""
    add_code_block(slide, "homework1/main.py · 时序切分与训练", train_code, 0.78, 1.38, 5.9, 4.78, size=12)
    add_code_block(slide, "homework1/backtest.py · 买卖信号与资金曲线", signal_code, 6.95, 1.38, 5.6, 4.78, size=12)
    add_rect(slide, 0.9, 6.32, 11.52, 0.48, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "代码要点：模型评估只使用 2024 年后 20% 测试集；最终回测模型使用全部 2024 年训练，2025 年数据只用于验证策略效果。", 1.12, 6.48, 11.1, 0.14, size=10, bold=True, color=COLORS["green"], align=PP_ALIGN.CENTER, margin=0)
    add_footer(slide, 7)

    # 8. ML evaluation
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "模型评估：误差水平接近", "2024 年测试集 RMSE / MAE，对比短期收益率预测精度")
    add_image_fit(slide, charts["metrics"], 0.78, 1.35, 5.75, 3.6)
    rows = [["模型", "RMSE", "MAE"]]
    for m in ["LinearRegression", "RandomForest", "LightGBM"]:
        rows.append([m, f"{metrics[m]['RMSE']:.6f}", f"{metrics[m]['MAE']:.6f}"])
    add_table(slide, rows, 7.0, 1.45, 5.05, 1.92, [2.15, 1.45, 1.45])
    add_rect(slide, 7.0, 3.84, 5.05, 1.42, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "解读", 7.28, 4.08, 0.8, 0.25, size=15, bold=True, color=COLORS["green"], margin=0)
    add_rich_lines(slide, ["三组模型 RMSE 均在 0.015 左右，约等于 1.5% 的日收益率误差。", "LinearRegression 的 RMSE 最低，RandomForest 的 MAE 最低；误差差距并不大。"], 7.28, 4.45, 4.25, 0.55, size=11, color=COLORS["muted"], line_spacing=0.9)
    add_footer(slide, 8)

    # 9. Prediction chart
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "预测效果：围绕零轴波动", "2025 年真实次日收益率与三组模型预测值对比")
    add_image_fit(slide, FIG_DIR / "prediction_vs_actual.png", 0.8, 1.27, 7.8, 3.95)
    # Add analysis metrics on the right
    add_rect(slide, 8.9, 1.27, 3.85, 3.95, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "关键发现", 9.2, 1.55, 2.0, 0.28, size=16, bold=True, color=COLORS["green"], margin=0)
    add_rich_lines(slide, [
        "预测值集中在 ±0.5% 范围内",
        "模型倾向于给出保守预测",
        "真实收益率波动更大（±3%）",
        "短期股价噪声难以消除",
    ], 9.2, 2.05, 3.2, 1.5, size=11, color=COLORS["muted"], line_spacing=1.1)
    add_rect(slide, 9.1, 3.75, 3.4, 1.2, fill=COLORS["mint"], line=COLORS["line"], radius=True)
    add_text(slide, "结论", 9.35, 3.95, 1.0, 0.22, size=12, bold=True, color=COLORS["green"], margin=0)
    add_rich_lines(slide, [
        "模型能捕捉部分趋势方向",
        "但收益率短期波动包含较强随机性",
    ], 9.35, 4.25, 2.9, 0.5, size=10, color=COLORS["ink"], line_spacing=0.9)
    add_footer(slide, 9)

    # 10. Backtest rules and result cards
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "量化回测：统一交易规则", "初始资金 100 万，交易手续费 0.03%，满仓买入 / 全部卖出")
    add_rect(slide, 0.78, 1.35, 11.78, 1.15, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    rules = [
        ("买入", "预测次日涨幅 > 0.5%，今日满仓买入"),
        ("卖出", "预测次日跌幅 > 0.5%，今日全部卖出"),
        ("持仓", "已持仓且无卖出信号则持有；空仓且无买入信号则观望"),
    ]
    for i, (name, desc) in enumerate(rules):
        x = 1.1 + i * 3.8
        add_rect(slide, x, 1.63, 0.52, 0.42, fill=[COLORS["green"], COLORS["red"], COLORS["navy"]][i], radius=True)
        add_text(slide, name, x + 0.09, 1.74, 0.36, 0.12, size=9, bold=True, color=COLORS["paper"], align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, desc, x + 0.68, 1.68, 2.85, 0.32, size=9, color=COLORS["ink"], margin=0)
    add_metric_card(slide, "LinearRegression", pct(backtests["LinearRegression"]["total_return"]), pct(backtests["LinearRegression"]["max_drawdown"]), pct(backtests["LinearRegression"]["win_rate"]), str(backtests["LinearRegression"]["trades"]), 0.88, 2.8, 3.72, 2.5)
    add_metric_card(slide, "RandomForest", pct(backtests["RandomForest"]["total_return"]), pct(backtests["RandomForest"]["max_drawdown"]), pct(backtests["RandomForest"]["win_rate"]), str(backtests["RandomForest"]["trades"]), 4.85, 2.8, 3.72, 2.5)
    add_metric_card(slide, "LightGBM", pct(backtests["LightGBM"]["total_return"]), pct(backtests["LightGBM"]["max_drawdown"]), pct(backtests["LightGBM"]["win_rate"]), str(backtests["LightGBM"]["trades"]), 8.82, 2.8, 3.72, 2.5)
    add_footer(slide, 10)

    # 11. Backtest curve
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "回测曲线：LightGBM 收益与回撤更优", "2025 年累计净值曲线，初始净值 = 1.0")
    add_image_fit(slide, FIG_DIR / "backtest_curves.png", 0.78, 1.28, 7.15, 4.55)
    add_image_fit(slide, charts["return_drawdown"], 8.2, 1.28, 4.65, 2.4)
    add_image_fit(slide, charts["buy_hold"], 8.2, 3.93, 4.65, 2.2)
    add_rect(slide, 0.9, 6.08, 11.62, 0.55, fill=COLORS["soft"], line=COLORS["line"], radius=True)
    add_text(slide, f"对照：2025 年买入持有收益约 {pct(buy_hold_2025)}；LightGBM 策略收益 {pct(backtests['LightGBM']['total_return'])}，说明交易信号在该样本期具备一定实盘价值。", 1.08, 6.27, 11.15, 0.18, size=11, bold=True, color=COLORS["green"], align=PP_ALIGN.CENTER, margin=0)
    add_footer(slide, 11)

    # 12. Strengths and limitations
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "模型优缺点与策略讨论", "机器学习指标接近，但交易表现差异明显")
    model_notes = [
        ("LinearRegression", "优点：简单、可解释、训练快", "缺点：难以捕捉非线性关系", "本次：RMSE最低，但回测小幅亏损"),
        ("RandomForest", "优点：可捕捉非线性，抗噪能力较强", "缺点：参数较多，交易信号偏弱", "本次：MAE最低，但累计收益最低"),
        ("LightGBM", "优点：提升树表达力强，风险控制较好", "缺点：交易次数更多，需要调参", "本次：唯一正收益，回撤最低"),
    ]
    for i, (model, good, bad, perf) in enumerate(model_notes):
        x = 0.82 + i * 4.05
        color = MODEL_COLORS[model]
        add_rect(slide, x, 1.42, 3.55, 2.25, fill=COLORS["paper"], line=COLORS["line"], radius=True)
        add_rect(slide, x, 1.42, 0.1, 2.25, fill=color)
        add_text(slide, model, x + 0.28, 1.68, 2.75, 0.25, size=16, bold=True, color=COLORS["ink"], margin=0)
        add_text(slide, good, x + 0.28, 2.15, 2.95, 0.22, size=11, color=COLORS["muted"], margin=0)
        add_text(slide, bad, x + 0.28, 2.55, 2.95, 0.22, size=11, color=COLORS["muted"], margin=0)
        add_text(slide, perf, x + 0.28, 3.0, 2.95, 0.22, size=11, bold=True, color=color, margin=0)
    # Merge strategy limitations and improvements into one section
    add_rect(slide, 0.82, 4.28, 11.7, 1.5, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "策略局限", 1.1, 4.54, 1.5, 0.25, size=16, bold=True, color=COLORS["red"], margin=0)
    add_rich_lines(slide, ["0.5% 阈值是统一规则，不一定最优", "满仓买入/全部卖出使资金曲线对信号质量敏感", "频繁交易会被手续费持续侵蚀"], 1.1, 4.92, 5.2, 0.65, size=11, color=COLORS["muted"], line_spacing=0.9)
    add_text(slide, "改进方向", 6.5, 4.54, 1.5, 0.25, size=16, bold=True, color=COLORS["green"], margin=0)
    add_rich_lines(slide, ["增加技术指标、行业指数等特征", "使用滚动训练或时间序列交叉验证", "优化阈值、止损规则与分批仓位管理"], 6.5, 4.92, 5.8, 0.65, size=11, color=COLORS["muted"], line_spacing=0.9)
    add_footer(slide, 12)

    # 13. Conclusion
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "结论", "本作业完成了从数据、特征、模型到回测的闭环验证")
    conclusions = [
        ("技术可行", "滚动窗口可形成短期预测框架，\n但单日收益率噪声较大。", COLORS["green"]),
        ("模型选择", "LightGBM 2025 年回测实现\n+10.40% 收益，回撤 5.94%。", COLORS["gold"]),
        ("投资应用", "需与阈值优化、仓位控制和\n风险管理结合，不能只看误差。", COLORS["navy"]),
    ]
    for i, (head, desc, color) in enumerate(conclusions):
        x = 0.75 + i * 4.08
        add_rect(slide, x, 1.55, 3.82, 2.8, fill=COLORS["paper"], line=COLORS["line"], radius=True)
        add_rect(slide, x, 1.55, 3.82, 0.08, fill=color)
        add_rect(slide, x + 0.22, 1.9, 0.5, 0.5, fill=color, radius=True)
        add_text(slide, str(i + 1), x + 0.35, 2.02, 0.2, 0.15, size=13, bold=True, color=COLORS["paper"], align=PP_ALIGN.CENTER, margin=0)
        add_text(slide, head, x + 0.9, 1.93, 2.5, 0.35, size=20, bold=True, color=COLORS["ink"], margin=0)
        add_rich_lines(slide, desc.split("\n"), x + 0.3, 2.7, 3.1, 1.2, size=11, color=COLORS["muted"], line_spacing=1.0)
    add_rect(slide, 0.75, 4.8, 11.85, 1.65, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_rect(slide, 0.75, 4.8, 0.08, 1.65, fill=COLORS["gold"])
    add_text(slide, "模型对比一览", 1.1, 5.05, 2.5, 0.28, size=15, bold=True, color=COLORS["ink"], margin=0)
    compare_rows = [
        ["模型", "累计收益", "最大回撤", "胜率", "交易次数"],
        ["LinearRegression", "-3.17%", "13.23%", "49.59%", "26"],
        ["RandomForest", "-7.95%", "12.41%", "47.52%", "25"],
        ["LightGBM", "+10.40%", "5.94%", "54.96%", "34"],
    ]
    add_table(slide, compare_rows, 1.1, 5.45, 8.5, 0.85, [2.1, 1.6, 1.6, 1.6, 1.6])
    add_rect(slide, 10.0, 5.05, 2.4, 1.2, fill=COLORS["mint"], line=COLORS["line"], radius=True)
    add_text(slide, "最佳策略", 10.25, 5.2, 1.8, 0.22, size=11, bold=True, color=COLORS["green"], margin=0)
    add_text(slide, "LightGBM", 10.25, 5.5, 1.8, 0.3, size=18, bold=True, color=COLORS["green"], margin=0)
    add_text(slide, "唯一正收益 + 最低回撤", 10.25, 5.88, 2.0, 0.18, size=9, color=COLORS["muted"], margin=0)
    add_rect(slide, 0.75, 6.7, 11.85, 0.42, fill=COLORS["soft"], line=COLORS["line"], radius=True)
    add_text(slide, "风险提示：回测结果不代表未来收益，本项目仅用于课程学习与方法演示，不构成投资建议。", 1.0, 6.82, 11.3, 0.18, size=10, bold=True, color=COLORS["muted"], align=PP_ALIGN.CENTER, margin=0)
    add_footer(slide, 13)

    # 14. Appendix
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["soft"])
    add_title(slide, "附录：项目结构与运行方式", "源码已模块化，可直接重新生成数据、图表与报告结果")
    add_rect(slide, 0.82, 1.38, 5.7, 4.45, fill=COLORS["soft"], line=COLORS["line"], radius=True)
    code_lines = [
        "Code/",
        "  common/data_fetcher.py   数据获取",
        "  homework1/config.py      参数配置",
        "  homework1/features.py    特征工程",
        "  homework1/models.py      模型训练与评估",
        "  homework1/backtest.py    量化回测",
        "  homework1/plots.py       结果可视化",
        "  homework1/main.py        主流程",
        "  outputs/homework1/       输出图表与PPT",
    ]
    add_rich_lines(slide, code_lines, 1.18, 1.72, 4.8, 3.6, size=11, color=COLORS["ink"], line_spacing=0.75)
    add_rect(slide, 7.08, 1.38, 5.15, 2.05, fill=COLORS["soft"], line=COLORS["line"], radius=True)
    add_text(slide, "运行命令", 7.4, 1.72, 1.6, 0.26, size=16, bold=True, color=COLORS["green"], margin=0)
    add_text(slide, "pip install -r requirements.txt\npython homework1/main.py", 7.4, 2.28, 4.35, 0.62, size=15, bold=True, color=COLORS["ink"], font="Consolas", margin=0)
    add_rect(slide, 7.08, 3.78, 5.15, 2.05, fill=COLORS["soft"], line=COLORS["line"], radius=True)
    add_text(slide, "依赖包", 7.4, 4.12, 1.4, 0.26, size=16, bold=True, color=COLORS["green"], margin=0)
    add_text(slide, "baostock · pandas · numpy · matplotlib · scikit-learn · lightgbm", 7.4, 4.72, 4.1, 0.45, size=13, color=COLORS["ink"], margin=0)
    add_footer(slide, 14)

    prs.save(PPTX_PATH)


def main():
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    result = compute_results()
    charts = create_support_charts(result)
    build_deck(result, charts)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
