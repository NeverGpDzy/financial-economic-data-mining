"""Build Homework 9 presentation from generated result files."""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "outputs" / "homework9"
PPT_PATH = OUTPUT_DIR / "202331060205_丁致宇_作业9_模型维护与风控对比实验.pptx"

COLORS = {
    "ink": "1F2937",
    "blue": "2563EB",
    "green": "059669",
    "gold": "D97706",
    "red": "DC2626",
    "paper": "F8FAFC",
    "line": "CBD5E1",
    "slate": "475569",
    "white": "FFFFFF",
}


def rgb(color: str) -> RGBColor:
    return RGBColor.from_string(color)


def add_bg(slide, color="F8FAFC"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def add_text(slide, text, x, y, w, h, size=16, bold=False, color="1F2937", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def add_title(slide, title, subtitle=None, dark=False):
    add_bg(slide, COLORS["ink"] if dark else COLORS["paper"])
    title_color = COLORS["white"] if dark else COLORS["ink"]
    sub_color = "D1D5DB" if dark else COLORS["slate"]
    add_text(slide, title, 0.55, 0.35, 12.0, 0.55, size=25, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.95, 11.9, 0.35, size=11.5, color=sub_color)
    if not dark:
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(COLORS["green"])
        bar.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent="059669"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("FFFFFF")
    shape.line.color.rgb = rgb(COLORS["line"])
    shape.line.width = Pt(0.8)
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.28, size=12.4, bold=True, color=accent)
    add_text(slide, body, x + 0.18, y + 0.58, w - 0.36, h - 0.68, size=10.5, color=COLORS["ink"])


def add_stat(slide, label, value, x, y, w, accent="059669"):
    add_text(slide, str(value), x, y, w, 0.52, size=25, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.55, w, 0.30, size=9.8, color=COLORS["slate"], align=PP_ALIGN.CENTER)


def add_fit_image(slide, path: Path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        draw_w = w
        draw_h = w / img_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * img_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))


def add_table(slide, df: pd.DataFrame, x, y, w, h, font_size=8.2):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(COLORS["green"])
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = "Microsoft YaHei"
        p.runs[0].font.size = Pt(font_size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = rgb(COLORS["white"])
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iloc[i, j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("FFFFFF" if i % 2 == 0 else "E2E8F0")
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.name = "Microsoft YaHei"
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.color.rgb = rgb(COLORS["ink"])


def pct(x: float) -> str:
    return f"{x * 100:.1f}%"


def load_results():
    summary = json.loads((OUTPUT_DIR / "summary.json").read_text(encoding="utf-8"))
    metrics = pd.read_csv(OUTPUT_DIR / "strategy_metrics.csv")
    windows = pd.read_csv(OUTPUT_DIR / "dynamic_windows.csv")
    return summary, metrics, windows


def build() -> None:
    summary, metrics, windows = load_results()
    calibration = summary["initial_calibration"]
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "作业9：模型维护与风控对比实验", "协整套利静态中枢、ADF动态风控与长期持有基准", dark=True)
    add_text(slide, "贵州茅台 vs 泸州老窖\n2015-2024 日度收盘价", 0.75, 2.0, 8.0, 0.88, size=27, bold=True, color=COLORS["white"])
    add_text(slide, "学生：丁致宇    学号：202331060205\n课程：金融与经济数据挖掘", 0.78, 5.82, 7.8, 0.65, size=13, color="D1D5DB")
    add_stat(slide, "B1窗口", summary["dynamic_window_count"], 8.5, 4.85, 1.35, COLORS["gold"])
    add_stat(slide, "B1通过", summary["dynamic_pass_window_count"], 10.05, 4.85, 1.35, COLORS["green"])
    add_stat(slide, "期初p值", f"{calibration['p_value']:.3f}", 11.55, 4.85, 1.2, COLORS["red"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "实验设计", "B1贴合教师指令，B2作为审查后的防未来函数对照。")
    add_card(slide, 0.75, 1.25, 3.55, 1.75, "数据与价差", "2015-2024日度收盘价\nspread = ln(茅台) - ln(老窖)", COLORS["blue"])
    add_card(slide, 4.85, 1.25, 3.55, 1.75, "方案A", "2015-2017一次性估计μ、σ\n2018-2024固定阈值不调整", COLORS["red"])
    add_card(slide, 8.95, 1.25, 3.35, 1.75, "方案B", "B1半年窗口内ADF\nB2最近三年滚动ADF", COLORS["green"])
    add_card(slide, 0.95, 4.15, 5.25, 1.35, "交易规则", "价差突破上轨做空价差，跌破下轨做多价差；等权多空，单腿50%。", COLORS["gold"])
    add_card(slide, 6.9, 4.15, 5.25, 1.35, "方案C", "2018首个交易日买入茅台或老窖并持有至2024年底，作为价值投资基准。", COLORS["blue"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "价格与价差", "竖线为2018年样本外回测起点。")
    add_fit_image(slide, OUTPUT_DIR / "price_and_spread.png", 0.75, 1.22, 11.8, 5.65)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "方案A：静态中枢", "固定期初μ和σ，后续中枢漂移不触发模型维护。")
    add_fit_image(slide, OUTPUT_DIR / "static_thresholds.png", 0.75, 1.22, 11.8, 5.65)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "方案B1：教师指令版动态风控", "红色淡区表示该半年窗口ADF未通过时空仓。")
    add_fit_image(slide, OUTPUT_DIR / "dynamic_thresholds.png", 0.75, 1.22, 11.8, 5.65)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "净值曲线对比", "配对套利与长期持有的收益来源不同。")
    add_fit_image(slide, OUTPUT_DIR / "strategy_nav_comparison.png", 0.75, 1.22, 11.8, 5.65)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "绩效指标", "累计收益、回撤和空仓时长共同反映风控差异。")
    show = metrics[["strategy", "cumulative_return", "max_drawdown", "trade_count", "cash_days", "risk_off_days"]].copy()
    show.columns = ["策略", "累计收益", "最大回撤", "调仓", "空仓天数", "风控空仓"]
    show["累计收益"] = show["累计收益"].map(pct)
    show["最大回撤"] = show["最大回撤"].map(pct)
    label_map = {
        "方案A：静态中枢无风控": "A 静态无风控",
        "方案B1：半年窗口ADF风控（教师指令版）": "B1 教师指令",
        "方案B2：滚动ADF风控（防未来函数）": "B2 防未来",
        "方案C1：贵州茅台长期持有": "C1 茅台持有",
        "方案C2：泸州老窖长期持有": "C2 老窖持有",
    }
    show["策略"] = show["策略"].map(label_map).fillna(show["策略"])
    add_table(slide, show, 0.65, 1.32, 12.1, 2.35, font_size=7.8)
    add_fit_image(slide, OUTPUT_DIR / "drawdown_comparison.png", 0.95, 4.15, 11.3, 2.45)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "半年ADF维护记录", "左图为B1教师指令版；B2防未来函数版保存在补充CSV和图表中。")
    add_fit_image(slide, OUTPUT_DIR / "dynamic_adf_pvalues.png", 0.75, 1.15, 7.15, 5.65)
    table = windows[["trade_window", "p_value", "can_trade", "mu_drift_vs_static"]].head(8).copy()
    table.columns = ["交易窗口", "p值", "交易", "中枢漂移"]
    table["p值"] = table["p值"].map(lambda x: f"{x:.3f}")
    table["交易"] = table["交易"].map(lambda x: "是" if x else "否")
    table["中枢漂移"] = table["中枢漂移"].map(lambda x: f"{x:.3f}")
    add_table(slide, table, 8.25, 1.35, 4.65, 4.5, font_size=7.0)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "结论", "模型维护不是提高收益的装饰项，而是统计套利能否继续交易的前提检查。")
    add_card(slide, 0.75, 1.25, 3.7, 1.75, "静态方案风险", "旧中枢长期固定，遇到估值结构变化时可能把漂移误判成套利信号。", COLORS["red"])
    add_card(slide, 4.85, 1.25, 3.7, 1.75, "动态风控价值", "B1贴合课堂指令，B2避免未来函数；两者都用ADF决定是否交易。", COLORS["green"])
    add_card(slide, 8.95, 1.25, 3.35, 1.75, "长期持有基准", "单边趋势收益强时可胜出，但风险来源和配对套利完全不同。", COLORS["blue"])
    add_card(slide, 0.9, 4.15, 11.25, 1.35, "最终判断", "协整套利策略必须持续监控ADF p值、中枢漂移、阈值稳定性和回撤路径；配对收益按等权多空总资金口径统计。", COLORS["gold"])

    prs.save(PPT_PATH)
    print(f"PPT已保存至: {PPT_PATH}")
    print(f"共{len(prs.slides)}页")


if __name__ == "__main__":
    build()
