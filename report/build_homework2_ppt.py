"""Generate a presentation for Homework 2 CAPM analysis."""

from __future__ import annotations

import json
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

OUT_DIR = ROOT / "outputs" / "homework2"
PPTX_PATH = OUT_DIR / "homework2_capm_presentation.pptx"

W, H = 13.333, 7.5

COLORS = {
    "ink": "1F2933",
    "muted": "667085",
    "paper": "FFFFFF",
    "soft": "F4F7F5",
    "green": "1B7F5C",
    "sage": "DDEBDD",
    "gold": "C89211",
    "red": "B43A3A",
    "navy": "25324D",
    "line": "D7DFE8",
}


def rgb(color: str) -> RGBColor:
    return RGBColor.from_string(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 16,
    bold: bool = False,
    color: str = COLORS["ink"],
    align=PP_ALIGN.LEFT,
    font: str = "Microsoft YaHei",
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rect(slide, x, y, w, h, *, fill=COLORS["paper"], line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def set_background(slide, color=COLORS["soft"]):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.62, 0.36, 8.8, 0.48, size=28, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.64, 0.92, 10.5, 0.32, size=11, color=COLORS["muted"])


def add_footer(slide, page):
    add_text(slide, "金融与经济数据挖掘 · 作业2 CAPM", 0.64, 7.15, 4.2, 0.22, size=8, color="98A2B3")
    add_text(slide, f"{page:02d}", 12.38, 7.15, 0.35, 0.22, size=8, color="98A2B3", align=PP_ALIGN.RIGHT)


def add_kpi(slide, label, value, x, y, w, h, color):
    add_rect(slide, x, y, w, h, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, value, x + 0.2, y + 0.2, w - 0.4, 0.42, size=26, bold=True, color=color)
    add_text(slide, label, x + 0.22, y + 0.78, w - 0.44, 0.24, size=10, color=COLORS["muted"])


def add_table(slide, rows, x, y, w, h, col_widths):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    for idx, width in enumerate(col_widths):
        table.columns[idx].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(COLORS["green"] if r == 0 else COLORS["paper"])
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(9)
                    run.font.bold = r == 0 or value == "泸州老窖"
                    run.font.color.rgb = rgb(COLORS["paper"] if r == 0 else COLORS["ink"])
    return shape


def pct(v: float) -> str:
    return f"{v:.2%}"


def load_results():
    capm = pd.read_csv(OUT_DIR / "capm_results.csv")
    curve = pd.read_csv(OUT_DIR / "backtest_curve.csv", parse_dates=["date"])
    summary = json.loads((OUT_DIR / "summary.json").read_text(encoding="utf-8"))
    return capm, curve, summary


def build_deck():
    capm, _curve, summary = load_results()
    metrics = summary["backtest_metrics"]
    extended = pd.DataFrame(summary.get("extended_backtest_results", []))

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["navy"])
    add_rect(slide, 0, 0, 0.25, H, fill=COLORS["gold"])
    add_text(slide, "《金融与经济数据挖掘》作业2", 0.78, 0.7, 5.0, 0.35, size=14, color="D9E6DD")
    add_text(slide, "CAPM 模型 Python 实现", 0.78, 1.45, 7.8, 0.7, size=38, bold=True, color=COLORS["paper"])
    add_text(slide, "沪深300作为市场基准 · p < 0.1 显著性筛选 · 2023-2024 买入持有回测", 0.8, 2.38, 8.6, 0.35, size=14, color="D9E6DD")
    add_kpi(slide, "显著最高 Alpha 股票", summary["selected_stock"], 0.9, 3.35, 2.75, 1.15, COLORS["gold"])
    add_kpi(slide, "日 Alpha", f"{summary['selected_alpha_daily']:.6f}", 3.95, 3.35, 2.15, 1.15, COLORS["green"])
    add_kpi(slide, "Alpha p 值", f"{summary['selected_alpha_pvalue']:.4f}", 6.4, 3.35, 2.0, 1.15, COLORS["red"])
    add_rect(slide, 9.05, 3.28, 3.25, 2.05, fill="31405F", line="536079", radius=True)
    add_text(slide, "核心结论", 9.32, 3.58, 1.6, 0.28, size=16, bold=True, color=COLORS["paper"])
    add_text(slide, f"历史显著 Alpha 没有在回测期兑现，泸州老窖 2023-2024 跑输沪深300 {abs(metrics['excess_total_return']):.2%}。", 9.32, 4.08, 2.65, 0.78, size=13, color="E6EDF3")
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "实验设计", "按作业要求从数据获取、CAPM 拟合到最高 Alpha 股票回测形成闭环")
    steps = [
        ("01", "数据", "Baostock 前复权日收盘价\n6只股票 + 沪深300"),
        ("02", "预处理", "计算日收益率\n1%/99% winsorize 缩尾"),
        ("03", "CAPM", "2020-2022 拟合\nRi-Rf = α + β(Rm-Rf)"),
        ("04", "筛选", "Alpha p < 0.1\n选择显著股票中 α 最高者"),
        ("05", "回测", "2023-2024 全仓买入持有\n对比沪深300"),
    ]
    for i, (num, head, body) in enumerate(steps):
        x = 0.78 + i * 2.5
        add_rect(slide, x, 1.65, 2.05, 3.25, fill=COLORS["paper"], line=COLORS["line"], radius=True)
        add_rect(slide, x + 0.22, 1.95, 0.52, 0.42, fill=COLORS["green"], radius=True)
        add_text(slide, num, x + 0.31, 2.05, 0.32, 0.14, size=9, bold=True, color=COLORS["paper"], align=PP_ALIGN.CENTER)
        add_text(slide, head, x + 0.22, 2.6, 1.5, 0.32, size=18, bold=True)
        add_text(slide, body, x + 0.24, 3.12, 1.55, 1.0, size=11, color=COLORS["muted"])
    add_rect(slide, 0.9, 5.55, 11.4, 0.72, fill=COLORS["sage"], line=COLORS["line"], radius=True)
    add_text(slide, "无风险日收益率固定采用 0.03 / 250 = 0.00012，符合本次作业要求。", 1.15, 5.78, 10.9, 0.22, size=13, bold=True, color=COLORS["green"], align=PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "CAPM 拟合结果", "泸州老窖是唯一在 10% 显著性水平下 Alpha 显著的股票")
    rows = [["股票", "日 Alpha", "年化 Alpha", "Beta", "Alpha p", "显著"]]
    for _, row in capm.iterrows():
        rows.append([
            row["stock"],
            f"{row['alpha_daily']:.6f}",
            pct(row["alpha_annualized"]),
            f"{row['beta']:.4f}",
            f"{row['alpha_pvalue']:.4f}",
            "是" if row["alpha_significant_10pct"] else "否",
        ])
    add_table(slide, rows, 0.72, 1.45, 8.25, 3.25, [1.45, 1.25, 1.35, 1.1, 1.15, 0.75])
    add_rect(slide, 9.35, 1.45, 3.05, 3.25, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "筛选结果", 9.62, 1.78, 1.5, 0.3, size=18, bold=True, color=COLORS["green"])
    add_text(slide, "泸州老窖", 9.62, 2.32, 2.1, 0.45, size=28, bold=True, color=COLORS["gold"])
    add_text(
        slide,
        f"α = {summary['selected_alpha_daily']:.6f}\np = {summary['selected_alpha_pvalue']:.4f}\nβ = {summary['selected_beta']:.4f}",
        9.65,
        3.04,
        2.1,
        0.9,
        size=16,
    )
    add_rect(slide, 0.92, 5.3, 11.3, 0.72, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "解释：泸州老窖的 Alpha 具有统计显著性，但 Beta 同时较高，意味着其对市场波动更加敏感。", 1.16, 5.53, 10.7, 0.22, size=13, bold=True, color=COLORS["navy"], align=PP_ALIGN.CENTER)
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "买入持有回测", "2023年初全仓买入泸州老窖，持有至2024年末，与沪深300对比")
    slide.shapes.add_picture(str(OUT_DIR / "backtest_curve.png"), Inches(0.72), Inches(1.42), Inches(7.15), Inches(4.2))
    add_kpi(slide, "泸州老窖累计收益率", pct(metrics["stock_total_return"]), 8.35, 1.55, 3.55, 1.05, COLORS["red"])
    add_kpi(slide, "沪深300累计收益率", pct(metrics["market_total_return"]), 8.35, 2.85, 3.55, 1.05, COLORS["green"])
    add_kpi(slide, "相对沪深300超额收益", pct(metrics["excess_total_return"]), 8.35, 4.15, 3.55, 1.05, COLORS["red"])
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "风险指标与反思", "高 Alpha 不等于低风险，更不等于未来一定有超额收益")
    rows = [
        ["指标", "泸州老窖", "沪深300"],
        ["期末资产", f"{metrics['stock_final_value']:,.2f}", f"{metrics['market_final_value']:,.2f}"],
        ["累计收益率", pct(metrics["stock_total_return"]), pct(metrics["market_total_return"])],
        ["最大回撤", pct(metrics["stock_max_drawdown"]), pct(metrics["market_max_drawdown"])],
        ["卡玛比率", f"{metrics['stock_calmar']:.4f}", f"{metrics['market_calmar']:.4f}"],
    ]
    add_table(slide, rows, 0.85, 1.55, 5.65, 2.6, [1.65, 2.0, 2.0])
    add_rect(slide, 7.0, 1.55, 5.2, 3.35, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "为什么理论 Alpha 没兑现？", 7.32, 1.9, 3.6, 0.32, size=19, bold=True, color=COLORS["green"])
    add_text(
        slide,
        "1. Alpha 只反映样本期平均超额收益，不保证未来持续性\n2. 2023-2024 市场风格切换，白酒板块估值回落\n3. 高 Beta 会放大市场下行期的回撤\n4. 单股全仓持有没有分散、止损和调仓\n5. CAPM 只解释市场因子，忽略行业和情绪因子",
        7.32,
        2.42,
        4.25,
        1.95,
        size=12,
        color=COLORS["ink"],
    )
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "补充实验：延长泸州老窖持有期", "保持选股结果不变，只把买入持有时间从 2024 年末继续拉长")
    slide.shapes.add_picture(str(OUT_DIR / "extended_horizon_comparison.png"), Inches(0.72), Inches(1.35), Inches(6.45), Inches(3.25))
    table_rows = [["持有期", "结束日", "泸州老窖", "相对沪深300"]]
    compact_names = ["2023-2024", "至2025年底", "至2026-05-18"]
    for idx, (_, row) in enumerate(extended.iterrows()):
        table_rows.append(
            [
                compact_names[idx] if idx < len(compact_names) else row["period"],
                row["end_date"],
                pct(row["stock_total_return"]),
                pct(row["excess_total_return"]),
            ]
        )
    add_table(slide, table_rows, 7.55, 1.42, 4.95, 1.95, [1.35, 1.25, 1.1, 1.25])
    add_rect(slide, 7.55, 3.75, 4.95, 1.75, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, "结果解释", 7.86, 4.02, 1.4, 0.26, size=17, bold=True, color=COLORS["green"])
    add_text(
        slide,
        "延长持有时间没有修复收益。沪深300后续逐步修复，但泸州老窖继续下行，说明历史 Alpha 并未持续，行业和个股风险超过了市场因子的解释范围。",
        7.86,
        4.47,
        4.15,
        0.72,
        size=12,
        color=COLORS["ink"],
    )
    add_rect(slide, 0.95, 5.85, 11.35, 0.58, fill=COLORS["sage"], line=COLORS["line"], radius=True)
    latest_row = extended.iloc[-1]
    add_text(
        slide,
        f"截至 {latest_row['end_date']}，泸州老窖累计收益 {latest_row['stock_total_return']:.2%}，相对沪深300落后 {abs(latest_row['excess_total_return']):.2%}。",
        1.16,
        6.05,
        10.85,
        0.18,
        size=12,
        bold=True,
        color=COLORS["navy"],
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_background(slide, COLORS["navy"])
    add_text(slide, "结论", 0.78, 0.82, 2.0, 0.55, size=34, bold=True, color=COLORS["paper"])
    add_rect(slide, 0.85, 1.85, 11.7, 3.65, fill="31405F", line="536079", radius=True)
    add_text(slide, f"在 p < 0.1 下，泸州老窖是唯一 Alpha 显著且 Alpha 最高的股票，日 Alpha 为 {summary['selected_alpha_daily']:.6f}，年化约 {summary['selected_alpha_annualized']:.2%}。", 1.18, 2.25, 10.95, 0.48, size=18, bold=True, color=COLORS["paper"])
    add_text(slide, f"但 2023-2024 回测中，泸州老窖累计收益率为 {metrics['stock_total_return']:.2%}，延长至 2026-05-18 后扩大到 {latest_row['stock_total_return']:.2%}。这说明 CAPM 历史 Alpha 是样本期统计结果，不具备天然的未来持续性。", 1.18, 3.12, 10.8, 0.72, size=18, color="E6EDF3")
    add_text(slide, "真实市场中获取 Alpha 需要结合基本面、市场状态、风险控制和动态再检验。", 1.18, 4.32, 10.8, 0.42, size=19, bold=True, color="F3D37A")
    add_footer(slide, 7)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    build_deck()
