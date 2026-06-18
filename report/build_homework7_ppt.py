"""Build Homework 7 presentation from generated result files."""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "outputs" / "homework7"
PPT_PATH = OUTPUT_DIR / "202331060205_丁致宇_作业7_CodeBuddy版平稳性检验.pptx"

COLORS = {
    "ink": "1F2937",
    "blue": "2563EB",
    "green": "059669",
    "gold": "D97706",
    "red": "DC2626",
    "paper": "F8FAFC",
    "line": "CBD5E1",
    "slate": "475569",
    "white": "FFFFFF",
}


def rgb(color: str) -> RGBColor:
    return RGBColor.from_string(color)


def add_bg(slide, color: str = "F8FAFC") -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def add_text(slide, text, x, y, w, h, size=16, bold=False, color="1F2937", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def add_title(slide, title, subtitle=None, dark=False):
    add_bg(slide, COLORS["ink"] if dark else COLORS["paper"])
    title_color = COLORS["white"] if dark else COLORS["ink"]
    sub_color = "D1D5DB" if dark else COLORS["slate"]
    add_text(slide, title, 0.55, 0.35, 11.8, 0.55, size=25, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.95, 11.8, 0.35, size=11.5, color=sub_color)
    if not dark:
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(COLORS["blue"])
        bar.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent="2563EB"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("FFFFFF")
    shape.line.color.rgb = rgb(COLORS["line"])
    shape.line.width = Pt(0.8)
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.28, size=12.5, bold=True, color=accent)
    add_text(slide, body, x + 0.18, y + 0.58, w - 0.36, h - 0.68, size=10.8, color=COLORS["ink"])


def add_stat(slide, label, value, x, y, w, accent="2563EB"):
    add_text(slide, str(value), x, y, w, 0.52, size=27, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.56, w, 0.30, size=10.3, color=COLORS["slate"], align=PP_ALIGN.CENTER)


def add_fit_image(slide, path: Path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        draw_w = w
        draw_h = w / img_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * img_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))


def add_table(slide, df: pd.DataFrame, x, y, w, h, font_size=8.6):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(COLORS["blue"])
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = "Microsoft YaHei"
        p.runs[0].font.size = Pt(font_size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = rgb(COLORS["white"])
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iloc[i, j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("FFFFFF" if i % 2 == 0 else "E2E8F0")
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.name = "Microsoft YaHei"
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.color.rgb = rgb(COLORS["ink"])


def load_results():
    summary = json.loads((OUTPUT_DIR / "summary.json").read_text(encoding="utf-8"))
    df = pd.read_csv(OUTPUT_DIR / "stationarity_summary.csv")
    return summary, df


def build() -> None:
    summary, df = load_results()
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "作业7：统计套利之平稳性检验", "CodeBuddy AI辅助编程版", dark=True)
    add_text(slide, "收盘价序列 vs 对数收益率序列的ADF单位根检验", 0.75, 2.0, 10.8, 0.62, size=27, bold=True, color=COLORS["white"])
    add_text(slide, "学生：丁致宇    学号：202331060205\n课程：金融与经济数据挖掘", 0.78, 5.82, 7.8, 0.65, size=13, color="D1D5DB")
    add_stat(slide, "价格平稳数量", f"{summary['price_stationary_count']}/5", 8.8, 4.85, 1.6, COLORS["gold"])
    add_stat(slide, "收益率平稳数量", f"{summary['return_stationary_count']}/5", 10.75, 4.85, 1.7, COLORS["green"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "实验设计", "五只股票统一使用2024年日度收盘价，分别检验价格和对数收益率。")
    add_card(slide, 0.75, 1.3, 3.6, 1.75, "标的股票", "中国石油、贵州茅台、兴蓉环境、招商银行、工商银行", COLORS["blue"])
    add_card(slide, 4.85, 1.3, 3.6, 1.75, "检验对象", "日度收盘价\n对数收益率 ln(Pt)-ln(Pt-1)", COLORS["green"])
    add_card(slide, 8.95, 1.3, 3.35, 1.75, "ADF判定", "p值 < 0.05 且 ADF统计量 < 5%临界值", COLORS["gold"])
    add_card(slide, 0.95, 4.1, 5.4, 1.35, "核心理论", "价格水平变量常含随机游走和趋势，容易非平稳；收益率是一阶差分后的相对变化，更接近围绕零均值波动。", COLORS["blue"])
    add_card(slide, 6.9, 4.1, 5.25, 1.35, "统计套利含义", "非平稳价格不能直接回归建模；若要使用价格，应先验证协整关系，否则优先使用收益率或平稳价差。", COLORS["green"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "ADF检验结果汇总", "p值和临界值共同决定平稳性结论。")
    show = df[["stock_name", "price_p_value", "price_stationarity", "return_p_value", "return_stationarity"]].copy()
    show.columns = ["股票", "价格p值", "价格结论", "收益率p值", "收益率结论"]
    show["价格p值"] = show["价格p值"].map(lambda x: f"{x:.4f}")
    show["收益率p值"] = show["收益率p值"].map(lambda x: f"{x:.4f}")
    add_table(slide, show, 0.8, 1.35, 11.7, 3.35, font_size=9.5)
    add_text(slide, "结论：收益率序列通常比价格序列更容易通过平稳性检验；若短样本价格通过检验，也只能说明该年度内呈现区间震荡，不代表长期价格平稳。", 0.95, 5.35, 11.1, 0.65, size=13, color=COLORS["ink"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "收盘价与收益率p值对比", "红色虚线为0.05显著性阈值。")
    add_fit_image(slide, OUTPUT_DIR / "adf_pvalue_comparison.png", 0.75, 1.35, 11.8, 5.45)

    for _, row in df.iterrows():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, f"{row['stock_name']}时序图", "上图为收盘价，下图为对数收益率。")
        add_fit_image(slide, OUTPUT_DIR / f"{row['code']}_{row['stock_name']}_price_return.png", 0.75, 1.35, 11.8, 5.45)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "理论思考题答案", "平稳性是统计套利信号可检验、可回测、可风控的基础。")
    add_card(slide, 0.75, 1.25, 3.7, 2.0, "价格为何非平稳", "价格持续吸收基本面、流动性和市场情绪信息，常带趋势或随机游走，均值方差不稳定。", COLORS["blue"])
    add_card(slide, 4.85, 1.25, 3.7, 2.0, "为何用收益率", "收益率去除价格水平趋势，更适合估计均值、方差、相关性和风险，能减少伪回归。", COLORS["green"])
    add_card(slide, 8.95, 1.25, 3.35, 2.0, "适合套利资产", "关系稳定、价差平稳、流动性好、交易成本低，并且基本面关系不易突然断裂。", COLORS["gold"])
    add_card(slide, 0.9, 4.35, 11.25, 1.15, "最终结论", "收益率平稳性检验验证了金融时序建模的基本前提；作业8进一步使用价格协整检验寻找可交易的平稳价差。", COLORS["blue"])

    prs.save(PPT_PATH)
    print(f"PPT已保存至: {PPT_PATH}")
    print(f"共{len(prs.slides)}页")


if __name__ == "__main__":
    build()
