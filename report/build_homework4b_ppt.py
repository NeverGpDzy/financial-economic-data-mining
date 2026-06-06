"""
Build the homework 4B presentation from freshly generated result files.
"""
import json
import os
import sys
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "outputs" / "homework4b"
PPT_PATH = OUTPUT_DIR / "202331060205_丁致宇_作业4B.pptx"

sys.path.insert(0, str(ROOT))


COLORS = {
    "ink": "16213E",
    "teal": "0B6E69",
    "mint": "43A047",
    "paper": "F7F9FA",
    "line": "D7DEE8",
    "slate": "4A5568",
    "accent": "F29E38",
    "red": "C2413B",
    "white": "FFFFFF",
}


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


def add_bg(slide, color="F7F9FA"):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    shape.z_order = 0


def add_text(slide, text, x, y, w, h, size=18, bold=False, color="16213E",
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = wrap
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def add_title(slide, title, subtitle=None, dark=False):
    if dark:
        add_bg(slide, COLORS["ink"])
        color = COLORS["white"]
        sub_color = "C7D2DE"
    else:
        add_bg(slide)
        color = COLORS["ink"]
        sub_color = COLORS["slate"]

    add_text(slide, title, 0.55, 0.35, 8.9, 0.55, size=25, bold=True, color=color)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.95, 11.0, 0.34, size=10.5, color=sub_color)
    if not dark:
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(COLORS["teal"])
        bar.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent="0B6E69"):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("FFFFFF")
    shape.line.color.rgb = rgb("D7DEE8")
    shape.line.width = Pt(0.8)
    add_text(slide, title, x + 0.2, y + 0.17, w - 0.4, 0.28, size=13, bold=True, color=accent)
    add_text(slide, body, x + 0.2, y + 0.58, w - 0.4, h - 0.75, size=11.5, color=COLORS["ink"])


def add_stat(slide, label, value, x, y, w, accent="0B6E69", label_color=None):
    value_text = str(value)
    value_size = 28 if len(value_text) >= 7 else 30
    add_text(slide, value_text, x, y, w, 0.52, size=value_size, bold=True,
             color=accent, align=PP_ALIGN.CENTER, wrap=False)
    add_text(slide, label, x, y + 0.55, w, 0.28, size=10.5,
             color=label_color or COLORS["slate"], align=PP_ALIGN.CENTER)


def add_table(slide, df, x, y, w, h, col_widths=None, font_size=9.5):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(COLORS["teal"])
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = "Microsoft YaHei"
        p.runs[0].font.size = Pt(font_size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = rgb(COLORS["white"])
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iloc[i, j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("FFFFFF" if i % 2 == 0 else "EDF3F6")
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.name = "Microsoft YaHei"
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.color.rgb = rgb(COLORS["ink"])
    return table


def add_fit_image(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        draw_w = w
        draw_h = w / img_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * img_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))


def fmt_float(value, digits=4):
    return f"{float(value):.{digits}f}"


def pct_from_summary(summary, key):
    return summary[key]


def load_results():
    summary = json.loads((OUTPUT_DIR / "summary.json").read_text(encoding="utf-8"))
    ir = pd.read_csv(OUTPUT_DIR / "ic_ir_results.csv")
    vif = pd.read_csv(OUTPUT_DIR / "vif_results.csv")
    ols = pd.read_csv(OUTPUT_DIR / "ols_results.csv")
    imp = pd.read_csv(OUTPUT_DIR / "feature_importance.csv")
    cv = pd.read_csv(OUTPUT_DIR / "cv_results.csv")
    bt = pd.read_csv(OUTPUT_DIR / "backtest_result.csv")
    trades = pd.read_csv(OUTPUT_DIR / "trade_log.csv")
    return summary, ir, vif, ols, imp, cv, bt, trades


def build():
    summary, ir, vif, ols, imp, cv, bt, trades = load_results()
    cumulative_excess = str(summary["超额收益(累计)"])
    excess_is_positive = not cumulative_excess.startswith("-")
    nav_subtitle = (
        "策略取得正收益并跑赢沪深300，但回撤控制仍需优化。"
        if excess_is_positive
        else "策略取得正收益，但样本外未跑赢沪深300，风险暴露仍需进一步约束。"
    )
    findings_subtitle = (
        "四因子中频短线流程可形成可验证的样本外收益，但风险暴露仍需进一步约束。"
        if excess_is_positive
        else "四因子中频短线流程取得正收益，但相对基准承压，风险暴露仍需进一步约束。"
    )
    conclusion_text = (
        f"核心结论：严格四因子口径和按日截面缩尾后，样本外累计收益{summary['累计收益率']}，"
        f"相对沪深300累计超额{summary['超额收益(累计)']}；结果低于AI初稿，但更符合题目要求和回测真实性。"
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, COLORS["ink"])
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.85), Inches(13.333), Inches(0.65)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = rgb(COLORS["teal"])
    slide.shapes[-1].line.fill.background()
    add_text(slide, "作业4B：中频短线量化全流程", 0.75, 1.35, 10.8, 0.72, size=34, bold=True, color=COLORS["white"])
    add_text(slide, "因子挖掘 → IC/共线质检 → LGBM赋权 → 1~2日持仓样本外回测", 0.8, 2.2, 9.8, 0.38, size=16, color="C7D2DE")
    add_text(slide, "学生：丁致宇    学号：202331060205\n课程：金融与经济数据挖掘    日期：2026年6月", 0.82, 5.95, 7.2, 0.7, size=12, color=COLORS["white"])
    add_stat(slide, "样本外累计收益", summary["累计收益率"], 9.25, 4.95, 1.95, COLORS["accent"], "C7D2DE")
    add_stat(slide, "最大回撤", summary["最大回撤"], 11.35, 4.95, 2.0, "F06A5D", "C7D2DE")

    # 2. Audit result
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "审查结论：旧版本存在口径偏差，已重算修正", "本版严格按作业4B要求复现四因子流程，样本外区间只用于回测。")
    add_card(slide, 0.65, 1.45, 3.7, 2.05, "发现的问题", "1. 原结果混入4个扩展因子\n2. IC/OLS覆盖了样本外年份\n3. 全样本缩尾存在预处理污染\n4. PPT写死旧指标，无法反映新结果", COLORS["red"])
    add_card(slide, 4.75, 1.45, 3.7, 2.05, "修复动作", "1. 因子改回Reversal/Liquidity/MoneyFlow/Value\n2. 质检仅使用2020-2023训练集\n3. 缩尾改为交易日截面3σ\n4. 回测修正成本、末日结算与基准起点", COLORS["teal"])
    add_card(slide, 8.85, 1.45, 3.45, 2.05, "最终口径", f"训练集：2020-2023\n样本外：{summary['回测区间']}\n股票池：沪深300成分股\n交易：每日Top5等权调仓", COLORS["mint"])
    add_text(slide, "本页用于说明对AI初稿的人工审查与修补，后续所有图表均来自修复后重跑结果。", 0.8, 4.65, 11.5, 0.45, size=15, bold=True, color=COLORS["ink"])

    # 3. Data design
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "数据与样本隔离设计", "本地Parquet数据：日频K线、daily_basic、5分钟线、财务指标、北向资金、沪深300指数。")
    add_stat(slide, "训练样本行数", "274,669", 0.8, 1.45, 2.2)
    add_stat(slide, "回测样本行数", "144,566", 3.2, 1.45, 2.2)
    add_stat(slide, "训练IC截面", "970", 5.6, 1.45, 2.2)
    add_stat(slide, "样本外交易日", summary["回测天数"], 8.0, 1.45, 2.2)
    add_stat(slide, "股票数量", "299", 10.4, 1.45, 2.2)
    add_card(slide, 0.8, 3.05, 5.8, 2.2, "时间切分", "训练集：2020-01-02 至 2023-12-29，用于因子质检、单因子OLS和LGBM训练。\n回测集：2024-01-02 至 2025-12-31，仅用于样本外打分、交易和绩效评价。")
    add_card(slide, 7.0, 3.05, 5.5, 2.2, "数据限制", "配套数据没有个股级 moneyflow_stock 的 buy_amount/sell_amount 字段，仅有市场级 moneyflow_hsgt。因此MoneyFlow采用个股成交额方向代理；去极值按交易日截面3σ执行。")

    # 4. Factors
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "四个中频短线因子定义", "只使用作业要求的4个因子，不再引入Momentum、Volatility、Turnover、VolumeChange。")
    factor_df = pd.DataFrame({
        "因子": ["Reversal", "Liquidity", "MoneyFlow", "Value"],
        "定义": ["-当日个股收益率", "|当日收益率| / 日成交额", "5日滚动签名成交额净流入 / 流通市值", "1 / PE_TTM"],
        "含义": ["短线反转", "交易摩擦/流动性", "个股资金流入流出代理", "低估值正向化"],
    })
    add_table(slide, factor_df, 0.75, 1.55, 11.8, 2.25, col_widths=[1.55, 5.0, 5.25], font_size=10)
    add_card(slide, 0.8, 4.55, 5.75, 1.55, "MoneyFlow代理口径", "签名成交额 = amount / 10 × sign(ret)，其中amount为千元、circ_mv为万元；先按股票滚动5日求和，再除以流通市值。")
    add_card(slide, 6.95, 4.55, 5.55, 1.55, "目标标签", "标签为个股次日超额收益 next_excess_ret，用于IC、OLS和LGBM训练；回测收益按持仓股票下一交易日ret结算。")

    # 5. IC/IR table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "训练集因子质检：IC/IR", "IC为当期因子与次日超额收益的日度截面Spearman相关；只在2020-2023训练集上计算。")
    ir_show = ir.copy()
    ir_show["IC均值"] = ir_show["IC均值"].map(lambda x: fmt_float(x, 4))
    ir_show["IR值"] = ir_show["IR值"].map(lambda x: fmt_float(x, 4))
    ir_show["t统计量"] = ir_show["t统计量"].map(lambda x: fmt_float(x, 2))
    ir_show = ir_show[["因子", "IC均值", "IR值", "t统计量", "有效判定", "样本数"]]
    add_table(slide, ir_show, 0.75, 1.35, 7.3, 2.15, col_widths=[1.45, 1.1, 1.1, 1.1, 1.15, 1.4], font_size=9.5)
    add_fit_image(slide, OUTPUT_DIR / "ic_ir_summary.png", 8.25, 1.05, 4.6, 3.3)
    add_card(slide, 0.9, 4.3, 3.6, 1.35, "有效因子", "Reversal、Liquidity、MoneyFlow达到|IC|≥0.02；Value为边缘因子，按课程规则保留进入后续模型。", COLORS["mint"])
    add_card(slide, 4.85, 4.3, 3.6, 1.35, "方向解释", "MoneyFlow的IC为负，说明在当前代理口径下，极端资金拥挤后的短期回撤更明显。", COLORS["accent"])
    add_card(slide, 8.8, 4.3, 3.45, 1.35, "稳定性", "Liquidity IR最高，为0.1997；Reversal与MoneyFlow也超过|0.1|稳定阈值。", COLORS["teal"])

    # 6. IC series
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "IC时间序列", "观察因子在训练期不同市场阶段下的预测能力波动。")
    add_fit_image(slide, OUTPUT_DIR / "ic_series.png", 0.6, 1.15, 12.0, 5.7)

    # 7. VIF and OLS
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "共线性与单因子OLS", "VIF用于排查因子冗余，OLS用于线性有效性量化；不因OLS失效而剔除非线性候选因子。")
    vif_show = vif.copy()
    vif_show["VIF"] = vif_show["VIF"].map(lambda x: fmt_float(x, 4))
    vif_show["R²"] = vif_show["R²"].map(lambda x: fmt_float(x, 4))
    add_table(slide, vif_show, 0.75, 1.35, 4.7, 2.35, col_widths=[1.55, 1.1, 1.1, 0.95], font_size=9.5)
    ols_show = ols[["因子", "β均值", "p值均值", "判定"]].copy()
    ols_show["β均值"] = ols_show["β均值"].map(lambda x: fmt_float(x, 6))
    ols_show["p值均值"] = ols_show["p值均值"].map(lambda x: fmt_float(x, 4))
    add_table(slide, ols_show, 6.0, 1.35, 6.25, 2.35, col_widths=[1.45, 1.45, 1.45, 1.9], font_size=9.3)
    add_card(slide, 0.9, 4.45, 5.2, 1.55, "VIF结论", "四个因子VIF均接近1，远低于VIF>5的高共线阈值，可全部纳入LGBM。", COLORS["mint"])
    add_card(slide, 6.7, 4.45, 5.35, 1.55, "OLS结论", "四个因子p值均高于0.05，线性维度待定；按作业规则全部进入LGBM做非线性二次筛选。", COLORS["accent"])

    # 8. LGBM
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "LGBM非线性赋权", "训练样本只来自2020-2023；特征为4个标准化因子，标签为次日超额收益。")
    cv_show = cv.copy()
    cv_show["train_mse"] = cv_show["train_mse"].map(lambda x: f"{x:.6f}")
    cv_show["val_mse"] = cv_show["val_mse"].map(lambda x: f"{x:.6f}")
    cv_show.columns = ["训练MSE", "验证MSE", "折数"]
    add_table(slide, cv_show[["折数", "训练MSE", "验证MSE"]], 0.75, 1.35, 4.6, 2.65, col_widths=[0.9, 1.8, 1.9], font_size=9.5)
    add_fit_image(slide, OUTPUT_DIR / "factor_weights.png", 5.65, 0.95, 6.85, 3.85)
    best = cv.loc[cv["val_mse"].idxmin()]
    add_card(slide, 0.9, 4.55, 4.4, 1.35, "交叉验证", f"5折扩展窗口时间序列CV，最佳为Fold {int(best['fold'])}，验证MSE={best['val_mse']:.6f}。", COLORS["teal"])
    add_card(slide, 5.9, 5.0, 5.9, 0.9, "重要性", "Liquidity权重最高，其次为Value、Reversal、MoneyFlow；模型自动学习非线性组合。", COLORS["mint"])

    # 9. Trading rules
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "日频截面选股与交易规则", "综合得分为LGBM预测的个股次日超额收益，得分越高越优先买入。")
    add_card(slide, 0.8, 1.35, 3.6, 2.0, "打分", "每日收盘后对沪深300股票池逐只预测次日超额收益，并按预测值从高到低排序。")
    add_card(slide, 4.75, 1.35, 3.6, 2.0, "选股", "每日选择Top5股票，等权持仓，每只股票目标权重20%。")
    add_card(slide, 8.7, 1.35, 3.6, 2.0, "成本", "初始建仓扣单边0.1%；后续调仓按卖出+买入实际换手部分扣成本。")
    add_card(slide, 0.8, 4.15, 5.65, 1.45, "时间对齐", "T日收盘产生信号并调仓；T+1日用持仓股票日收益结算。最后一个回测日只结算持仓，不再新开无效仓位。", COLORS["teal"])
    add_card(slide, 6.9, 4.15, 5.4, 1.45, "交易日志", f"输出每日调仓记录，共{len(trades)}条；包含持仓、买入、卖出、换手率和成本。", COLORS["accent"])

    # 10. Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "样本外回测核心指标", "回测区间：2024-01-02 至 2025-12-31，初始资金100万元，单边交易成本0.1%。")
    stats = [
        ("累计收益率", summary["累计收益率"], COLORS["teal"]),
        ("年化收益率", summary["年化收益率"], COLORS["teal"]),
        ("最大回撤", summary["最大回撤"], COLORS["red"]),
        ("日胜率", summary["日胜率"], COLORS["accent"]),
        ("累计超额", summary["超额收益(累计)"], COLORS["mint"]),
        ("夏普比率", str(summary["夏普比率"]), COLORS["ink"]),
    ]
    for idx, (label, value, color) in enumerate(stats):
        add_stat(slide, label, value, 0.8 + idx * 2.0, 1.35, 1.7, color)
    add_fit_image(slide, OUTPUT_DIR / "metrics_table.png", 2.0, 3.05, 9.3, 3.15)

    # 11. NAV
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "策略净值 vs 沪深300", nav_subtitle)
    add_fit_image(slide, OUTPUT_DIR / "nav_vs_market.png", 0.55, 1.05, 12.2, 5.85)

    # 12. Drawdown
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "回撤分析", f"最大回撤为{summary['最大回撤']}，主要来自高波动阶段的日频满仓换仓风险。")
    add_fit_image(slide, OUTPUT_DIR / "drawdown.png", 0.75, 1.1, 11.7, 5.45)

    # 13. Findings
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "研究结论", findings_subtitle)
    add_card(slide, 0.8, 1.35, 3.7, 2.0, "因子层", "Reversal、Liquidity、MoneyFlow在训练集IC层面有效；Value为边缘因子但对LGBM非线性组合有贡献。", COLORS["mint"])
    add_card(slide, 4.85, 1.35, 3.7, 2.0, "模型层", "LGBM可处理OLS难以刻画的非线性关系；特征重要性显示流动性和估值是主要分裂收益来源。", COLORS["teal"])
    add_card(slide, 8.9, 1.35, 3.5, 2.0, "交易层", "每日Top5换仓带来较高信号响应速度，也带来较高交易成本和回撤风险。", COLORS["accent"])
    add_text(slide, "下一步优化方向：加入仓位上限、行业中性、止损/止盈规则，并对MoneyFlow代理口径做稳健性检验。", 0.9, 4.7, 11.3, 0.6, size=16, bold=True, color=COLORS["ink"])

    # 14. AI reflection
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "AI辅助编程审查反思", "AI能快速生成流程，但量化作业必须进行数据口径、时序逻辑和结果真实性审查。")
    add_card(slide, 0.8, 1.3, 5.6, 2.1, "AI初稿风险", "AI倾向为了提升结果加入额外因子；也容易将全样本统计结果写入PPT，造成样本外信息污染。")
    add_card(slide, 6.9, 1.3, 5.5, 2.1, "人工修复价值", "人工审查将因子口径、训练/回测隔离、交易成本、最后交易日处理逐项校正，使结果可解释、可复现。")
    add_card(slide, 0.8, 4.25, 5.6, 1.4, "保留限制", "MoneyFlow受限于本地数据字段，只能使用个股级代理；这会影响经济解释强度。", COLORS["accent"])
    add_card(slide, 6.9, 4.25, 5.5, 1.4, "提交材料", "PPT、AI代码审查记录、AI交互记录、回测图表/指标表格均已按修复后结果更新。", COLORS["mint"])

    # 15. Submission
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "提交文件清单", "命名遵循：学号 + 姓名 + 作业4B。")
    files_df = pd.DataFrame({
        "类别": ["PPT", "代码", "数据", "图表", "指标表", "AI记录"],
        "文件/目录": [
            "202331060205_丁致宇_作业4B.pptx / .pdf",
            "homework4b/",
            "data/homework4b/train_data.csv、test_data.csv、full_data.csv、数据说明.md",
            "outputs/homework4b/*.png",
            "summary.json、IC/VIF/OLS/feature_importance/backtest_result",
            "AI代码审核记录.md、AI交互记录.md、AI交互记录.docx、AI代码审查与修复表.docx",
        ],
        "说明": ["演示文稿", "全流程源码", "清洗后训练/回测/全量数据", "净值、回撤、IC、重要性", "可复核数据结果", "AI辅助与人工修复过程"],
    })
    add_table(slide, files_df, 0.75, 1.45, 11.8, 3.65, col_widths=[1.25, 6.25, 4.3], font_size=8.7)
    add_text(slide, conclusion_text, 0.95, 5.45, 11.2, 0.65, size=15, bold=True, color=COLORS["ink"])

    prs.save(PPT_PATH)
    print(f"PPT已保存至: {PPT_PATH}")
    print(f"共{len(prs.slides)}页")


if __name__ == "__main__":
    build()
