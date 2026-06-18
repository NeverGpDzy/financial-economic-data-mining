"""Build Homework 8 presentation from generated result files."""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "outputs" / "homework8"
PPT_PATH = OUTPUT_DIR / "202331060205_丁致宇_作业8_协整套利投资.pptx"

COLORS = {
    "ink": "1F2937",
    "blue": "2563EB",
    "green": "059669",
    "gold": "D97706",
    "red": "DC2626",
    "paper": "F8FAFC",
    "line": "CBD5E1",
    "slate": "475569",
    "white": "FFFFFF",
}


def rgb(color: str) -> RGBColor:
    return RGBColor.from_string(color)


def add_bg(slide, color="F8FAFC"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def add_text(slide, text, x, y, w, h, size=16, bold=False, color="1F2937", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def add_title(slide, title, subtitle=None, dark=False):
    add_bg(slide, COLORS["ink"] if dark else COLORS["paper"])
    title_color = COLORS["white"] if dark else COLORS["ink"]
    sub_color = "D1D5DB" if dark else COLORS["slate"]
    add_text(slide, title, 0.55, 0.35, 11.9, 0.55, size=25, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.95, 11.9, 0.35, size=11.5, color=sub_color)
    if not dark:
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(COLORS["green"])
        bar.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent="059669"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("FFFFFF")
    shape.line.color.rgb = rgb(COLORS["line"])
    shape.line.width = Pt(0.8)
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.28, size=12.5, bold=True, color=accent)
    add_text(slide, body, x + 0.18, y + 0.58, w - 0.36, h - 0.68, size=10.7, color=COLORS["ink"])


def add_stat(slide, label, value, x, y, w, accent="059669"):
    add_text(slide, str(value), x, y, w, 0.52, size=26, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.56, w, 0.30, size=10.2, color=COLORS["slate"], align=PP_ALIGN.CENTER)


def add_fit_image(slide, path: Path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        draw_w = w
        draw_h = w / img_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * img_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))


def add_table(slide, df: pd.DataFrame, x, y, w, h, font_size=8.0):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(COLORS["green"])
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = "Microsoft YaHei"
        p.runs[0].font.size = Pt(font_size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = rgb(COLORS["white"])
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iloc[i, j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("FFFFFF" if i % 2 == 0 else "E2E8F0")
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.name = "Microsoft YaHei"
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.color.rgb = rgb(COLORS["ink"])


def load_results():
    summary = json.loads((OUTPUT_DIR / "summary.json").read_text(encoding="utf-8"))
    pairs = pd.read_csv(OUTPUT_DIR / "cointegration_pairs.csv")
    return summary, pairs


def build() -> None:
    summary, pairs = load_results()
    best = summary["best_pair"]
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "作业8：协整检验与配对交易模型", "统计套利之协整套利投资", dark=True)
    add_text(slide, "EG两步法协整检验 → 最优配对 → 价差z-score择时", 0.75, 2.0, 10.8, 0.62, size=27, bold=True, color=COLORS["white"])
    add_text(slide, "学生：丁致宇    学号：202331060205\n课程：金融与经济数据挖掘", 0.78, 5.82, 7.8, 0.65, size=13, color="D1D5DB")
    add_stat(slide, "检验配对", summary["pair_count"], 8.7, 4.85, 1.35, COLORS["gold"])
    add_stat(slide, "显著协整", summary["significant_pair_count"], 10.3, 4.85, 1.35, COLORS["green"])
    add_stat(slide, "最优p值", f"{best['p_value']:.4f}", 11.85, 4.85, 1.1, COLORS["red"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "实验流程", "从价格非平稳前提出发，寻找平稳价差并生成交易信号。")
    add_card(slide, 0.75, 1.25, 2.65, 1.7, "1. 数据获取", "获取标的池2015-2018日度收盘价，并按交易日期对齐。", COLORS["blue"])
    add_card(slide, 3.7, 1.25, 2.65, 1.7, "2. EG检验", "OLS拟合Y=α+βX+ε，对残差ε做ADF检验。", COLORS["green"])
    add_card(slide, 6.65, 1.25, 2.65, 1.7, "3. 最优配对", "按残差ADF p值升序排序，选显著性最强的一组。", COLORS["gold"])
    add_card(slide, 9.6, 1.25, 2.65, 1.7, "4. 交易信号", "价差z-score超过±2开仓，回归0附近平仓观察。", COLORS["red"])
    add_card(slide, 0.9, 4.25, 5.45, 1.25, "标的池", "中国石油、贵州茅台、泸州老窖、兴蓉环境、招商银行、工商银行、长江电力、上证50指数。", COLORS["green"])
    add_card(slide, 6.9, 4.25, 5.25, 1.25, "判断标准", "残差p值<0.05且ADF统计量小于5%临界值，说明价差平稳，配对具备统计套利基础。", COLORS["blue"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "协整检验结果Top10", "按残差ADF p值排序，p值越小，协整显著性越强。")
    show = pairs.head(10)[["rank", "y_asset", "x_asset", "beta", "r_squared", "p_value", "conclusion"]].copy()
    show.columns = ["排名", "Y资产", "X资产", "beta", "R²", "p值", "结论"]
    show["Y资产"] = show["Y资产"].str.split("（").str[0]
    show["X资产"] = show["X资产"].str.split("（").str[0]
    show["beta"] = show["beta"].map(lambda x: f"{x:.3f}")
    show["R²"] = show["R²"].map(lambda x: f"{x:.3f}")
    show["p值"] = show["p值"].map(lambda x: f"{x:.4f}")
    add_table(slide, show, 0.55, 1.25, 12.25, 4.6, font_size=7.8)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "协整p值排序图", "红色虚线为0.05显著性阈值。")
    add_fit_image(slide, OUTPUT_DIR / "cointegration_pvalues_top12.png", 0.75, 1.35, 11.8, 5.45)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "最优配对标准化价格走势", f"{best['y_asset']} vs {best['x_asset']}")
    add_fit_image(slide, OUTPUT_DIR / "best_pair_standardized_prices.png", 0.75, 1.35, 11.8, 5.45)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "线性拟合与价差平稳性", "价差即OLS残差，是配对交易的核心监控对象。")
    add_stat(slide, "beta", f"{best['ols_beta']:.3f}", 0.8, 1.25, 1.55, COLORS["green"])
    add_stat(slide, "R²", f"{best['ols_r_squared']:.3f}", 2.55, 1.25, 1.55, COLORS["blue"])
    add_stat(slide, "ADF p值", f"{best['p_value']:.4f}", 4.3, 1.25, 1.55, COLORS["red"])
    add_fit_image(slide, OUTPUT_DIR / "best_pair_spread.png", 0.75, 2.35, 11.8, 4.3)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "z-score交易信号", "z-score超过±2时价差偏离较大，回归0附近时平仓观察。")
    add_fit_image(slide, OUTPUT_DIR / "best_pair_zscore_signals.png", 0.75, 1.35, 11.8, 5.45)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "深度问答", "协整关系需要持续维护，不能把历史显著性当作未来保证。")
    add_card(slide, 0.75, 1.25, 3.7, 1.65, "为何出现配对", "共同基本面、同行业/同指数、相似风险因子暴露，使线性价差可能均值回复。", COLORS["green"])
    add_card(slide, 4.85, 1.25, 3.7, 1.65, "未来一定重现吗", "不一定。协整是历史统计关系，基本面、流动性、监管和资金结构变化都可能破坏关系。", COLORS["red"])
    add_card(slide, 8.95, 1.25, 3.35, 1.65, "失效因素", "重大事件、行业分化、交易拥挤、极端行情、样本过拟合、成本上升。", COLORS["gold"])
    add_card(slide, 0.9, 4.15, 11.25, 1.25, "维护模型", "滚动重估协整关系和对冲系数，监控残差ADF p值、z-score分布、价差均值漂移、交易成本与止损规则；关系恶化时暂停或更换配对。", COLORS["blue"])

    prs.save(PPT_PATH)
    print(f"PPT已保存至: {PPT_PATH}")
    print(f"共{len(prs.slides)}页")


if __name__ == "__main__":
    build()
