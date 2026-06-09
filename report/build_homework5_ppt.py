"""Build Homework 5 presentation from generated result files."""

from __future__ import annotations

import json
import sys
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "outputs" / "homework5"
PPT_PATH = OUTPUT_DIR / "202331060205_丁致宇_作业5.pptx"

COLORS = {
    "ink": "1F2937",
    "teal": "0F766E",
    "mint": "14B8A6",
    "gold": "D97706",
    "paper": "F8FAFC",
    "line": "CBD5E1",
    "slate": "475569",
    "red": "DC2626",
    "white": "FFFFFF",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_bg(slide, color: str = "F8FAFC") -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 16,
    bold: bool = False,
    color: str = "1F2937",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def add_title(slide, title: str, subtitle: str | None = None, dark: bool = False) -> None:
    add_bg(slide, COLORS["ink"] if dark else COLORS["paper"])
    title_color = COLORS["white"] if dark else COLORS["ink"]
    sub_color = "D1D5DB" if dark else COLORS["slate"]
    add_text(slide, title, 0.55, 0.35, 11.5, 0.55, size=26, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.95, 11.6, 0.35, size=11.5, color=sub_color)
    if not dark:
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(COLORS["teal"])
        bar.line.fill.background()


def add_card(slide, x, y, w, h, title, body, accent="0F766E"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("FFFFFF")
    shape.line.color.rgb = rgb(COLORS["line"])
    shape.line.width = Pt(0.8)
    add_text(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.3, size=13, bold=True, color=accent)
    add_text(slide, body, x + 0.18, y + 0.58, w - 0.36, h - 0.72, size=11.2, color=COLORS["ink"])


def add_stat(slide, label, value, x, y, w, accent="0F766E"):
    add_text(slide, str(value), x, y, w, 0.52, size=28, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.56, w, 0.28, size=10.5, color=COLORS["slate"], align=PP_ALIGN.CENTER)


def add_fit_image(slide, path: Path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        draw_w = w
        draw_h = w / img_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * img_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))


def add_table(slide, df: pd.DataFrame, x, y, w, h, font_size=8.8):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(COLORS["teal"])
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = "Microsoft YaHei"
        p.runs[0].font.size = Pt(font_size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = rgb(COLORS["white"])
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iloc[i, j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("FFFFFF" if i % 2 == 0 else "E2E8F0")
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.name = "Microsoft YaHei"
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.color.rgb = rgb(COLORS["ink"])
    return table


def pct(value: float) -> str:
    return f"{float(value):.2%}"


def load_results():
    summary = json.loads((OUTPUT_DIR / "summary.json").read_text(encoding="utf-8"))
    primary = summary["period_results"]["2022-2024"]
    robust = summary["period_results"]["2022-2023"]
    top20 = pd.read_csv(OUTPUT_DIR / "top20_train_alpha_2022_2024.csv")
    groups = pd.read_csv(OUTPUT_DIR / "group_persistence_2022_2024.csv")
    return summary, primary, robust, top20, groups


def build() -> None:
    summary, primary, robust, top20, groups = load_results()
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "作业5：Alpha的持续性", "AI辅助的上证50指数股CAPM实证分析", dark=True)
    add_text(slide, "历史高Alpha股票，未来还能保持高Alpha吗？", 0.75, 2.0, 9.6, 0.65, size=28, bold=True, color=COLORS["white"])
    add_text(slide, "学生：丁致宇    学号：202331060205\n课程：金融与经济数据挖掘", 0.78, 5.85, 7.8, 0.65, size=13, color="D1D5DB")
    add_stat(slide, "主口径Top20重合度", pct(primary["overlap_ratio"]), 9.2, 4.85, 1.8, COLORS["gold"])
    add_stat(slide, "排名相关", f"{primary['spearman_corr']:.3f}", 11.05, 4.85, 1.55, COLORS["mint"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "实验设计", "按CAPM分别估计训练期与未来期Alpha，再检验Top20的持续性。")
    add_card(slide, 0.75, 1.35, 3.5, 2.05, "训练期", "2019-01-01 至 2021-12-31\n逐股回归 CAPM\n按年化Alpha降序排序", COLORS["teal"])
    add_card(slide, 4.85, 1.35, 3.5, 2.05, "主检验期", "2022-01-01 至 2024-12-31\n对应AI指令第4条\n计算未来Alpha与排名", COLORS["gold"])
    add_card(slide, 8.95, 1.35, 3.35, 2.05, "稳健性", "2022-01-01 至 2023-12-31\n对应实验数据说明\n用于消除口径争议", COLORS["mint"])
    add_card(slide, 0.9, 4.25, 5.45, 1.35, "市场基准", f"{summary['data_audit']['market_source']}，无风险利率为年化1.5%。")
    add_card(slide, 6.9, 4.25, 5.35, 1.35, "持续性指标", "Top20重合度、Alpha均值变化、Spearman排名相关、按历史Alpha分组后的未来Alpha。")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "主结果概览：2022-2024", "历史高Alpha组合在未来期整体出现明显衰减。")
    add_stat(slide, "重合度", pct(primary["overlap_ratio"]), 0.85, 1.45, 1.7, COLORS["teal"])
    add_stat(slide, "重合股票数", f"{primary['overlap_count']}/{primary['top_n']}", 2.95, 1.45, 1.7, COLORS["teal"])
    add_stat(slide, "训练期Top20均值", pct(primary["train_top_alpha_mean"]), 5.05, 1.45, 1.9, COLORS["gold"])
    add_stat(slide, "后期同组均值", pct(primary["train_top_future_alpha_mean"]), 7.45, 1.45, 1.9, COLORS["gold"])
    add_stat(slide, "Alpha变化", pct(primary["alpha_mean_change"]), 9.85, 1.45, 1.8, COLORS["red"])
    add_card(slide, 0.9, 3.55, 5.4, 1.5, "核心判断", f"主口径下，前期Top20平均Alpha从 {pct(primary['train_top_alpha_mean'])} 下降到 {pct(primary['train_top_future_alpha_mean'])}，Alpha持续性较弱。", COLORS["red"])
    add_card(slide, 6.85, 3.55, 5.15, 1.5, "排名关系", f"全样本前后Alpha Spearman相关为 {primary['spearman_corr']:.4f}，历史排序对未来排序解释力有限。", COLORS["teal"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Alpha散点：训练期 vs 未来期", "橙色点为训练期Top20，虚线为Alpha不变的45度线。")
    add_fit_image(slide, OUTPUT_DIR / "alpha_scatter_2022_2024.png", 0.75, 1.1, 11.8, 5.7)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "训练期Top20前后Alpha对比", "大部分历史高Alpha股票在未来期Alpha回落。")
    add_fit_image(slide, OUTPUT_DIR / "top20_alpha_compare_2022_2024.png", 0.75, 1.05, 11.8, 5.85)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Top20重合度", "重合度越高，说明历史高Alpha组合越能延续到未来。")
    add_fit_image(slide, OUTPUT_DIR / "overlap_2022_2024.png", 0.65, 1.0, 5.7, 5.65)
    add_fit_image(slide, OUTPUT_DIR / "group_persistence_2022_2024.png", 6.4, 1.05, 6.25, 5.55)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "稳健性：2022-2023口径", "题目数据说明中的样本外区间给出类似结论。")
    add_stat(slide, "稳健口径重合度", pct(robust["overlap_ratio"]), 0.85, 1.55, 2.2, COLORS["teal"])
    add_stat(slide, "同组后期Alpha", pct(robust["train_top_future_alpha_mean"]), 3.55, 1.55, 2.2, COLORS["gold"])
    add_stat(slide, "Alpha变化", pct(robust["alpha_mean_change"]), 6.25, 1.55, 2.2, COLORS["red"])
    add_stat(slide, "Spearman", f"{robust['spearman_corr']:.3f}", 8.95, 1.55, 1.8, COLORS["mint"])
    add_card(slide, 0.95, 3.8, 5.4, 1.5, "结论一致", "无论采用2022-2024还是2022-2023，历史Top20组合的未来Alpha均明显低于训练期。")
    add_card(slide, 6.9, 3.8, 5.25, 1.5, "评分口径兼容", "报告和结果表同时保留两个区间，课堂展示时可按老师解释选择主结果。", COLORS["mint"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "训练期Top20明细", "按2019-2021年化Alpha从高到低排序。")
    show = top20.head(10)[["code", "alpha_rank_train", "alpha_rank_test", "alpha_annual_train", "alpha_annual_test", "alpha_change"]].copy()
    show.columns = ["代码", "训练排名", "后期排名", "训练Alpha", "后期Alpha", "变化"]
    for col in ["训练Alpha", "后期Alpha", "变化"]:
        show[col] = show[col].map(pct)
    add_table(slide, show, 0.75, 1.3, 11.8, 4.3, font_size=8.5)
    add_text(slide, "注：本页展示前10只，完整Top20见 outputs/homework5/top20_train_alpha_2022_2024.csv。", 0.9, 6.05, 11.2, 0.35, size=11, color=COLORS["slate"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "金融解释与应对方案", "Alpha衰减不是代码错误，而是市场竞争和时变结构下的常见现象。")
    add_card(slide, 0.75, 1.25, 3.6, 2.1, "原因1：均值回归", "短期超额收益可能来自误定价修复，修复后高Alpha自然下降。", COLORS["gold"])
    add_card(slide, 4.85, 1.25, 3.6, 2.1, "原因2：风格轮动", "不同时期市场偏好的行业、估值和风险暴露不同，历史强势风格可能切换。", COLORS["teal"])
    add_card(slide, 8.95, 1.25, 3.35, 2.1, "原因3：交易拥挤", "高Alpha信号被更多资金使用后，收益被提前兑现，未来超额收益下降。", COLORS["red"])
    add_card(slide, 0.9, 4.35, 5.45, 1.4, "如何避免", "滚动窗口更新Alpha；做行业/风格中性；加入估值、质量和风险控制；必须保留样本外检验。", COLORS["mint"])
    add_card(slide, 6.85, 4.35, 5.3, 1.4, "展示结论", "历史高Alpha可作为候选信号，但不能单独作为未来选股依据，需要与多因子和风控体系结合。", COLORS["teal"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "AI辅助代码反思", "AI提升了生成速度，但口径审查和金融解释仍必须人工完成。")
    add_card(slide, 0.8, 1.35, 5.6, 1.85, "AI生成有用之处", "快速搭建数据读取、CAPM回归、排序分组、图表输出和报告框架。")
    add_card(slide, 6.9, 1.35, 5.5, 1.85, "人工关键修改", "补齐市场基准、处理题目区间矛盾、显式加入无风险利率、输出稳健性结果。", COLORS["gold"])
    add_card(slide, 0.8, 4.1, 5.6, 1.45, "结果解读", "不能只看训练期Alpha；真正重要的是未来期表现、重合度和稳定性。", COLORS["teal"])
    add_card(slide, 6.9, 4.1, 5.5, 1.45, "提交材料", "源码、结果表、图表、Markdown报告、PPT均在仓库内可复现生成。", COLORS["mint"])

    prs.save(PPT_PATH)
    print(f"PPT已保存至: {PPT_PATH}")
    print(f"共{len(prs.slides)}页")


if __name__ == "__main__":
    build()

