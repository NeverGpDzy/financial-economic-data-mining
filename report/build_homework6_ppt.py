"""Build Homework 6 presentation from generated result files."""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_DIR = ROOT / "outputs" / "homework6"
PPT_PATH = OUTPUT_DIR / "202331060205_丁致宇_作业6.pptx"

COLORS = {
    "ink": "1D2B23",
    "forest": "1B7F5A",
    "moss": "8FB339",
    "gold": "D19A32",
    "brick": "A23E48",
    "paper": "F7F8F3",
    "line": "D6D8CF",
    "slate": "475569",
    "white": "FFFFFF",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_bg(slide, color: str = "F7F8F3") -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()


def add_text(slide, text, x, y, w, h, size=15, bold=False, color="1D2B23", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    for paragraph in tf.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return box


def add_title(slide, title, subtitle=None, dark=False):
    add_bg(slide, COLORS["ink"] if dark else COLORS["paper"])
    title_color = COLORS["white"] if dark else COLORS["ink"]
    sub_color = "DDE5D6" if dark else COLORS["slate"]
    add_text(slide, title, 0.55, 0.35, 11.6, 0.58, size=25, bold=True, color=title_color)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.95, 11.7, 0.35, size=11.5, color=sub_color)


def add_card(slide, x, y, w, h, title, body, accent="1B7F5A"):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("FFFFFF")
    shape.line.color.rgb = rgb(COLORS["line"])
    shape.line.width = Pt(0.8)
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.36, 0.32, size=12.5, bold=True, color=accent)
    add_text(slide, body, x + 0.18, y + 0.58, w - 0.36, h - 0.72, size=10.8, color=COLORS["ink"])


def add_stat(slide, label, value, x, y, w, accent="1B7F5A", label_color=None):
    if label_color is None:
        label_color = COLORS["slate"]
    add_text(slide, value, x, y, w, 0.52, size=27, bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.56, w, 0.32, size=10.2, color=label_color, align=PP_ALIGN.CENTER)


def add_fit_image(slide, path: Path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio >= box_ratio:
        draw_w = w
        draw_h = w / img_ratio
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        draw_h = h
        draw_w = h * img_ratio
        draw_x = x + (w - draw_w) / 2
        draw_y = y
    slide.shapes.add_picture(str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h))


def add_table(slide, df: pd.DataFrame, x, y, w, h, font_size=8.4):
    rows, cols = df.shape
    table = slide.shapes.add_table(rows + 1, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(COLORS["forest"])
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.runs[0].font.name = "Microsoft YaHei"
        p.runs[0].font.size = Pt(font_size)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = rgb(COLORS["white"])
    for i in range(rows):
        for j in range(cols):
            cell = table.cell(i + 1, j)
            cell.text = str(df.iloc[i, j])
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("FFFFFF" if i % 2 == 0 else "EEF1E8")
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.name = "Microsoft YaHei"
            p.runs[0].font.size = Pt(font_size)
            p.runs[0].font.color.rgb = rgb(COLORS["ink"])
    return table


def pct(value) -> str:
    return f"{float(value):.2%}"


def load_results():
    summary = json.loads((OUTPUT_DIR / "summary.json").read_text(encoding="utf-8"))
    ic = pd.read_csv(OUTPUT_DIR / "factor_quality_summary.csv")
    imp = pd.read_csv(OUTPUT_DIR / "feature_importance.csv")
    fcff = pd.read_csv(OUTPUT_DIR / "fcff_group_summary.csv")
    metrics = pd.read_csv(OUTPUT_DIR / "price_metrics.csv")
    annual = pd.read_csv(OUTPUT_DIR / "price_annual_returns.csv")
    return summary, ic, imp, fcff, metrics, annual


def build() -> None:
    summary, ic, imp, fcff, metrics, annual = load_results()
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "作业6B：基于FCFF的价值投资", "传统规则选股 + LGBM多因子价值模型双实证", dark=True)
    add_text(slide, "巴菲特·芒格“三好价值投资”的量化落地", 0.75, 2.0, 9.2, 0.65, size=27, bold=True, color=COLORS["white"])
    add_text(slide, "学生：丁致宇    学号：202331060205\n课程：金融与经济数据挖掘", 0.78, 5.75, 7.8, 0.75, size=13, color="DDE5D6")
    a_metrics = metrics[metrics["group"].eq("A")]
    best = a_metrics.sort_values("年化收益率", ascending=False).iloc[0]
    add_stat(slide, "A组最佳年化收益", pct(best["年化收益率"]), 8.8, 4.75, 2.0, COLORS["gold"], "DDE5D6")
    add_stat(slide, "覆盖股票数", str(summary["data_audit"]["stock_count"]), 10.95, 4.75, 1.5, COLORS["moss"], "DDE5D6")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "实验目标", "将“三好”理念拆成财务因子，并同时检验传统规则与机器学习打分。")
    add_card(slide, 0.8, 1.35, 3.4, 2.0, "好商业模式", "净利润增速、毛利率、净利率、轻资产、低负债，衡量企业赚钱方式和资产效率。", COLORS["forest"])
    add_card(slide, 4.95, 1.35, 3.4, 2.0, "经济护城河", "ROE与费用率，观察公司是否能长期维持资本回报和经营效率。", COLORS["moss"])
    add_card(slide, 9.1, 1.35, 3.25, 2.0, "长期现金流", "利润含金量、分红率、股息率、FCFF收益率，聚焦真实可分配现金。", COLORS["gold"])
    add_card(slide, 0.9, 4.35, 5.35, 1.35, "方案A", "固定财务阈值硬筛选，透明稳定但阈值僵硬。")
    add_card(slide, 6.95, 4.35, 5.25, 1.35, "方案B", "LGBM非线性赋权，能捕捉交互关系但必须控制过拟合。", COLORS["brick"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "数据口径", "按教师提供文件可复现执行，并显式说明与题目年份的差异。")
    audit = summary["data_audit"]
    add_stat(slide, "实际行情区间", f"{audit['trade_date_min']} ~ {audit['trade_date_max']}", 0.8, 1.4, 4.0, COLORS["forest"])
    add_stat(slide, "年度财务截面", f"{audit['panel_year_min']}-{audit['panel_year_max']}", 5.05, 1.4, 2.4, COLORS["gold"])
    add_stat(slide, "年度面板行数", str(audit["panel_rows"]), 7.9, 1.4, 1.8, COLORS["moss"])
    add_stat(slide, "股票数", str(audit["stock_count"]), 10.1, 1.4, 1.6, COLORS["brick"])
    add_card(slide, 0.95, 3.85, 5.4, 1.55, "防未来泄露", "未来1年FCFF标签用于训练，因此模型训练标签截至2016；2017只作为2018持仓打分起点。", COLORS["brick"])
    add_card(slide, 6.85, 3.85, 5.3, 1.55, "样本外范围", "价格回测覆盖2018-2024；未来3年FCFF标签因数据截止2024，只能验证到2021年截面。", COLORS["forest"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "整体流程", "数据加载 → 因子质检 → LGBM打分 → 双策略分组 → FCFF与股价回测。")
    steps = [
        ("1 数据", "年度财务、分红、复权行情、沪深300"),
        ("2 因子", "F1-F11正向化、缩尾、Z标准化"),
        ("3 质检", "IC/IR、VIF、单因子OLS"),
        ("4 模型", "树深≤3的LGBM非线性赋权"),
        ("5 回测", "A/B/C分层，年度调仓等权持有"),
    ]
    for i, (title, body) in enumerate(steps):
        add_card(slide, 0.65 + i * 2.52, 2.05, 2.15, 2.25, title, body, [COLORS["forest"], COLORS["moss"], COLORS["gold"], COLORS["brick"], COLORS["forest"]][i])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "因子质检结果", "IC均值、IR和VIF共同决定因子是否适合进入模型。")
    add_fit_image(slide, OUTPUT_DIR / "factor_ic_ir.png", 0.55, 1.1, 7.0, 5.7)
    top_ic = ic.sort_values("IC均值", ascending=False).head(5)[["因子", "IC均值", "IR", "VIF"]].copy()
    top_ic["IC均值"] = top_ic["IC均值"].map(lambda v: f"{v:.4f}")
    top_ic["IR"] = top_ic["IR"].map(lambda v: f"{v:.3f}")
    top_ic["VIF"] = top_ic["VIF"].map(lambda v: "" if pd.isna(v) else f"{v:.2f}")
    add_table(slide, top_ic, 7.85, 1.45, 4.85, 3.35, font_size=7.4)
    add_text(slide, "表格展示IC均值最高的前5个因子。", 8.0, 5.15, 4.4, 0.35, size=10.5, color=COLORS["slate"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "LGBM非线性赋权", "小样本下限制树深和叶子数，重点看样本外分层而不是训练拟合。")
    add_fit_image(slide, OUTPUT_DIR / "feature_importance.png", 0.75, 1.05, 6.55, 5.65)
    cv = pd.read_csv(OUTPUT_DIR / "lgbm_cv_results.csv")
    if not cv.empty:
        add_card(slide, 7.85, 1.45, 4.6, 1.35, "时间序列CV", f"验证折数：{len(cv)}\n平均验证MSE：{cv['val_mse'].mean():.4f}", COLORS["forest"])
    add_card(slide, 7.85, 3.35, 4.6, 1.65, "控参", "max_depth=3，num_leaves=7，min_data_in_leaf=5，L2正则控制过拟合。", COLORS["brick"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "两套选股方案", "同一年度截面，同一A/B/C分组，后续FCFF和股价回测使用完全一致的分组股票。")
    add_card(slide, 0.9, 1.3, 5.3, 2.1, "方案A：固定阈值传统规则", "净利润正增长、净利率为正、轻资产、低负债、ROE质量、费用控制、现金流质量、分红与FCFF为正。", COLORS["forest"])
    add_card(slide, 6.95, 1.3, 5.25, 2.1, "方案B：LGBM综合打分", "使用过VIF质检的标准化因子预测未来1年FCFF增速，分数越高表示中长期价值预期越好。", COLORS["brick"])
    add_card(slide, 0.9, 4.35, 11.3, 1.25, "统一回测规则", "年末截面打分，下一年等权持有；初始资金100万，单边手续费0.1%，基准为沪深300。", COLORS["gold"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "FCFF分层回测", "检验A组真实未来3年FCFF年化增速是否高于B、C组。")
    add_fit_image(slide, OUTPUT_DIR / "fcff_group_backtest.png", 0.75, 1.0, 11.8, 5.85)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "股价收益回测", "A组组合与沪深300基准的样本外净值对比。")
    add_fit_image(slide, OUTPUT_DIR / "price_nav_comparison.png", 0.75, 1.0, 11.8, 5.85)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "年度收益与核心指标", "年度收益用于观察策略是否只依赖单一年份。")
    add_fit_image(slide, OUTPUT_DIR / "annual_returns_A_group.png", 0.55, 1.05, 6.35, 5.75)
    add_fit_image(slide, OUTPUT_DIR / "price_metrics_table.png", 6.95, 1.5, 5.65, 2.2)
    add_fit_image(slide, OUTPUT_DIR / "top_group_radar.png", 8.05, 3.85, 3.8, 3.0)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "结果分析与思考题", "价值模型的关键不是短期涨跌，而是企业现金流质量能否转化为长期回报。")
    add_card(slide, 0.75, 1.25, 3.65, 2.05, "护城河优先", "高ROE和费用控制代表企业能长期保留竞争优势，区别于只看价格惯性的短线策略。", COLORS["forest"])
    add_card(slide, 4.85, 1.25, 3.65, 2.05, "现金流优先", "净利润容易受会计确认影响，FCFF更接近股东真实可分配资金。", COLORS["gold"])
    add_card(slide, 8.95, 1.25, 3.4, 2.05, "模型差异", "传统规则稳定可解释，LGBM灵活但更依赖样本外检验和正则控制。", COLORS["brick"])
    add_card(slide, 0.9, 4.35, 5.35, 1.35, "与4B区别", "本模型是年度、基本面、低频持有；4B是日频行情和流动性因子，持仓周期短。")
    add_card(slide, 6.9, 4.35, 5.25, 1.35, "最终结论", "价值因子可以构成候选股票池，但必须结合样本外验证、行业暴露和现金流质量判断。", COLORS["moss"])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "总结 & Q&A", "AI负责生成和迭代代码，人工负责口径审查、无泄露约束和金融解释。", dark=True)
    add_text(slide, "交付材料：源码、因子质检表、单因子回归、特征重要度、双策略入选股票、FCFF与股价回测、报告、AI记录和PPT。", 0.9, 2.15, 11.2, 0.85, size=21, bold=True, color=COLORS["white"])
    add_text(slide, "关键审查点：实际数据范围、未来标签、年度调仓、手续费、沪深300基准。", 0.95, 4.55, 10.8, 0.6, size=14, color="DDE5D6")

    prs.save(PPT_PATH)
    print(f"PPT已保存至: {PPT_PATH}")
    print(f"共{len(prs.slides)}页")


if __name__ == "__main__":
    build()
