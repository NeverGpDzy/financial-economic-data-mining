"""Generate a presentation for Homework 3: Two-Factor Model."""

from __future__ import annotations

import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

import json
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT_DIR = ROOT / "outputs" / "homework3"
PPTX_PATH = OUT_DIR / "homework3_presentation.pptx"

W, H = 13.333, 7.5

COLORS = {
    "ink": "202124",
    "muted": "667085",
    "soft": "F5F7FA",
    "paper": "FFFFFF",
    "green": "1B7F5C",
    "red": "B5403C",
    "navy": "22304D",
    "line": "D8DEE8",
    "blue": "2F6FED",
    "orange": "D97904",
    "purple": "6E56CF",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_text(slide, text, x, y, w, h, *, size=18, color=COLORS["ink"],
             bold=False, font="Microsoft YaHei", align=PP_ALIGN.LEFT, margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_lines(slide, lines, x, y, w, h, *, size=14, color=COLORS["ink"], line_spacing=1.15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(4 * line_spacing)
    return box


def add_rect(slide, x, y, w, h, *, fill=COLORS["paper"], line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_title(slide, title, subtitle=None, *, inverse=False):
    title_color = COLORS["paper"] if inverse else COLORS["ink"]
    sub_color = "D7E2D8" if inverse else COLORS["muted"]
    add_text(slide, title, 0.65, 0.36, 8.6, 0.55, size=26, bold=True, color=title_color, margin=0)
    if subtitle:
        add_text(slide, subtitle, 0.67, 0.95, 9.2, 0.32, size=10, color=sub_color, margin=0)


def add_footer(slide, page, inverse=False):
    c = "D7E2D8" if inverse else "98A2B3"
    add_text(slide, "金融与经济数据挖掘 · 作业3", 0.65, 7.16, 4.2, 0.22, size=8, color=c, margin=0)
    add_text(slide, f"{page:02d}", 12.35, 7.16, 0.35, 0.22, size=8, color=c, align=PP_ALIGN.RIGHT, margin=0)


def add_table(slide, headers, rows, x, y, w, row_h=0.35, col_widths=None, header_color=COLORS["navy"]):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(row_h * n_rows))
    table = table_shape.table

    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(header_color)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = "Microsoft YaHei"
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = rgb(COLORS["paper"])

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(COLORS["soft"])
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(9)
                    run.font.color.rgb = rgb(COLORS["ink"])
    return table_shape


def build_ppt():
    summary_path = OUT_DIR / "summary.json"
    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    comp = pd.read_csv(OUT_DIR / "model_comparison.csv")
    metrics = summary["backtest_metrics"]
    growth = summary["growth_rates"]

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # ========== Slide 1: Title ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W, H, fill=COLORS["navy"])
    add_rect(slide, 0, 5.8, W, 0.04, fill=COLORS["blue"])
    add_text(slide, "作业3：二因子模型", 0.8, 1.8, 11, 0.8, size=40, bold=True, color=COLORS["paper"], margin=0)
    add_text(slide, "CAPM + PE_TTM 估值因子  |  PEG 策略回测", 0.85, 2.8, 11, 0.5, size=18, color="A0AEC0", margin=0)
    add_text(slide, "金融与经济数据挖掘", 0.85, 4.5, 6, 0.35, size=13, color="7B8794", margin=0)

    # ========== Slide 2: 实验概述 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "实验概述", "模型构建、回归对比与策略回测")
    add_footer(slide, 2)

    lines = [
        "一、实验目的",
        "  · 掌握 CAPM 单因子与二因子模型的构建、回归与对比方法",
        "  · 理解估值因子 PE_TTM 对股票收益的解释作用",
        "  · 基于二因子模型的估值逻辑，使用 PEG 策略进行样本外回测",
        "",
        "二、核心流程",
        "  数据获取 → 异常值处理 → CAPM 拟合 → 二因子模型拟合 → 模型对比 → PEG 回测",
    ]
    add_rich_lines(slide, lines, 0.65, 1.5, 12, 5.5, size=14)

    # ========== Slide 3: 数据说明 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "数据说明", "baostock 数据接口 | 6 只 A 股 + 上证指数")
    add_footer(slide, 3)

    headers = ["股票", "代码", "行业", "利润增长率"]
    rows = []
    stock_industry = {"贵州茅台": "白酒", "中国石油": "能源", "五粮液": "白酒",
                      "泸州老窖": "白酒", "招商银行": "金融", "美的集团": "家电"}
    stock_codes = {"贵州茅台": "sh.600519", "中国石油": "sh.601857", "五粮液": "sz.000858",
                   "泸州老窖": "sz.000568", "招商银行": "sh.600036", "美的集团": "sz.000333"}
    for name in stock_codes:
        rows.append([name, stock_codes[name], stock_industry[name], f"{growth[name]:.2%}"])
    add_table(slide, headers, rows, 0.65, 1.6, 8.5, col_widths=[1.8, 1.8, 1.2, 1.5])

    info_lines = [
        "时间区间：",
        f"  训练期：{summary['training_period']}",
        f"  回测期：{summary['backtest_period']}",
        "",
        "无风险利率：1.5%/年（日化）",
        "市场基准：上证指数 (sh.000001)",
    ]
    add_rich_lines(slide, info_lines, 9.5, 1.6, 3.5, 4, size=12, color=COLORS["muted"])

    # ========== Slide 4: 数据预处理 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "数据预处理", "收益率计算 | Winsorize 缩尾 | PE_TTM 归一化")
    add_footer(slide, 4)

    lines = [
        "1. 收益率计算",
        "   return = close.pct_change()，取每日价格变化百分比",
        "",
        "2. 异常值处理 — Winsorize 缩尾 (limits=[0.01, 0.01])",
        "   · 对收益率进行缩尾，截断 1%/99% 分位数以外的极端值",
        "   · 对 PE_TTM 进行缩尾，处理无穷大、负值等极端值",
        "",
        "3. PE_TTM 归一化 — 最大最小归一化",
        "   PE_norm = (PE_winsorized - PE_min) / (PE_max - PE_min)",
        "   · PE_min/PE_max 基于训练期（2018-2022）计算",
        "   · 回测期使用相同的训练期极值进行归一化",
        "",
        "4. 无风险利率：Rf = 1.5%/年，日化 = (1+0.015)^(1/250) - 1",
    ]
    add_rich_lines(slide, lines, 0.65, 1.5, 12, 5.5, size=13)

    # ========== Slide 5: CAPM 模型结果 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "CAPM 模型拟合结果", "Ri - Rf = alpha + beta x (Rm - Rf)")
    add_footer(slide, 5)

    headers = ["股票", "alpha(日)", "alpha(年化)", "beta(mkt)", "alpha p值", "R^2"]
    rows = []
    for _, r in comp.iterrows():
        rows.append([
            r["股票"],
            f"{r['CAPM_α']:.6f}",
            f"{r['CAPM_α'] * 250:.2%}",
            f"{r['CAPM_β(mkt)']:.4f}",
            f"{r['CAPM_α_p值']:.4f}",
            f"{r['CAPM_R²']:.4f}",
        ])
    add_table(slide, headers, rows, 0.65, 1.5, 12, col_widths=[1.8, 1.6, 1.6, 1.6, 1.6, 1.2])

    sig_line = "* alpha p值 < 0.10 为显著"
    add_text(slide, sig_line, 0.65, 5.8, 6, 0.3, size=10, color=COLORS["muted"])

    # ========== Slide 6: 二因子模型结果 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "二因子模型拟合结果", "Ri - Rf = alpha + beta1 x (Rm - Rf) + beta2 x PE_norm")
    add_footer(slide, 6)

    headers = ["股票", "alpha(日)", "beta(mkt)", "beta(pe)", "alpha p值", "beta(pe) p值", "R^2"]
    rows = []
    for _, r in comp.iterrows():
        rows.append([
            r["股票"],
            f"{r['二因子_α']:.6f}",
            f"{r['二因子_β(mkt)']:.4f}",
            f"{r['二因子_β(pe)']:.4f}",
            f"{r['二因子_α_p值']:.4f}",
            f"{r['二因子_β(pe)_p值']:.4f}",
            f"{r['二因子_R²']:.4f}",
        ])
    add_table(slide, headers, rows, 0.65, 1.5, 12, col_widths=[1.6, 1.5, 1.4, 1.4, 1.5, 1.5, 1.0])

    sig_line = "* beta(pe) p值 < 0.10 为显著   |   招商银行 p=0.023 (5%显著), 美的集团 p=0.077 (10%显著)"
    add_text(slide, sig_line, 0.65, 5.8, 10, 0.3, size=10, color=COLORS["muted"])

    # ========== Slide 7: 模型对比分析 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "模型对比分析", "CAPM vs 二因子模型")
    add_footer(slide, 7)

    # Insert comparison plot
    plot_path = OUT_DIR / "model_comparison.png"
    if plot_path.exists():
        slide.shapes.add_picture(str(plot_path), Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.0))

    lines = [
        "关键发现：",
        "  1. 加入 PE_TTM 后 R^2 提升极小（<0.3%），说明 PE_TTM 对日度收益增量解释力有限",
        "  2. 白酒三股 alpha 绝对值减少，PE_TTM 部分吸收了超额收益异象",
        "  3. 仅招商银行和美的集团的 beta(pe) 统计显著",
    ]
    add_rich_lines(slide, lines, 0.65, 5.5, 12, 1.8, size=12)

    # ========== Slide 8: 因子暴露与行业分析 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "因子暴露与行业差异分析")
    add_footer(slide, 8)

    headers = ["股票", "行业", "beta(pe)", "p值", "显著性", "分析"]
    rows = []
    analysis = {
        "贵州茅台": "白酒：估值因子不显著，受情绪驱动",
        "中国石油": "能源：低beta，受油价等外部因素驱动",
        "五粮液": "白酒：估值因子不显著，高波动",
        "泸州老窖": "白酒：beta(pe)接近显著，估值有一定解释力",
        "招商银行": "金融：beta(pe)最显著，估值波动有规律",
        "美的集团": "家电：beta(pe)接近显著，估值与收益有关联",
    }
    for _, r in comp.iterrows():
        name = r["股票"]
        sig = "显著" if r["二因子_β(pe)_p值"] < 0.10 else "不显著"
        rows.append([name, stock_industry[name], f"{r['二因子_β(pe)']:.4f}",
                      f"{r['二因子_β(pe)_p值']:.4f}", sig, analysis[name]])
    add_table(slide, headers, rows, 0.65, 1.5, 12, row_h=0.45, col_widths=[1.4, 0.9, 1.2, 1.0, 0.9, 5.0])

    # ========== Slide 9: PEG 策略说明 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "PEG 策略回测", "基于二因子模型估值逻辑的择时策略")
    add_footer(slide, 9)

    lines = [
        "PEG 计算：",
        "  PEG = 归一化后 PE_TTM / 五年利润平均增长率（2018-2022 算术平均）",
        "",
        "交易规则：",
        "  ① PEG < 0.8  → 全仓买入（多只满足选 alpha 最优者）",
        "  ② PEG > 1.5  → 全仓卖出",
        "  ③ 0.8 <= PEG <= 1.5  → 维持当前持仓",
        "",
        f"交易成本：手续费 {summary['commission']:.2%}/笔 + 滑点 {summary['slippage']:.1%}/笔",
        f"初始资金：{summary['initial_capital']:,.0f} 元",
        f"回测期间：{summary['backtest_period']}",
    ]
    add_rich_lines(slide, lines, 0.65, 1.5, 12, 5.5, size=14)

    # ========== Slide 10: 回测结果 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "回测结果", "PEG 策略 vs 上证指数")
    add_footer(slide, 10)

    # KPI cards
    def kpi(label, value, x, color):
        add_rect(slide, x, 1.6, 2.8, 1.4, fill=COLORS["paper"], line=COLORS["line"], radius=True)
        add_text(slide, value, x + 0.15, 1.75, 2.5, 0.55, size=28, bold=True, color=color, margin=0)
        add_text(slide, label, x + 0.15, 2.45, 2.5, 0.35, size=11, color=COLORS["muted"], margin=0)

    kpi("策略累计收益", f"{metrics['strategy_total_return']:.2%}", 0.65, COLORS["blue"])
    kpi("市场累计收益", f"{metrics['market_total_return']:.2%}", 3.7, COLORS["orange"])
    kpi("策略最大回撤", f"{metrics['strategy_max_drawdown']:.2%}", 6.75, COLORS["red"])
    kpi("策略卡玛比率", f"{metrics['strategy_calmar']:.4f}", 9.8, COLORS["green"])

    # Trade info
    lines = [
        f"交易次数：{metrics['trade_count']} 次",
        f"持仓天数：{metrics['holding_days']} 天（{metrics['holding_pct']:.1f}%）",
        f"策略年化收益：{metrics['strategy_annual_return']:.2%}  |  市场年化收益：{metrics['market_annual_return']:.2%}",
    ]
    add_rich_lines(slide, lines, 0.65, 3.3, 12, 1.5, size=13)

    # Insert backtest curve
    curve_path = OUT_DIR / "backtest_curve.png"
    if curve_path.exists():
        slide.shapes.add_picture(str(curve_path), Inches(0.4), Inches(4.6), Inches(12.5), Inches(2.6))

    # ========== Slide 11: 回撤分析 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "回撤分析")
    add_footer(slide, 11)

    dd_path = OUT_DIR / "drawdown.png"
    if dd_path.exists():
        slide.shapes.add_picture(str(dd_path), Inches(0.4), Inches(1.5), Inches(12.5), Inches(4.5))

    lines = [
        f"策略最大回撤：{metrics['strategy_max_drawdown']:.2%}  |  市场最大回撤：{metrics['market_max_drawdown']:.2%}",
        "最大回撤出现在 2024 年初，市场整体下行期间",
    ]
    add_rich_lines(slide, lines, 0.65, 6.2, 12, 1, size=12, color=COLORS["muted"])

    # ========== Slide 12: 交易记录 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "交易记录")
    add_footer(slide, 12)

    trade_log = pd.read_csv(OUT_DIR / "trade_log.csv")
    headers = ["日期", "操作", "标的", "组合价值", "备注"]
    trade_rows = []
    for _, t in trade_log.iterrows():
        action_map = {"buy": "买入", "sell": "卖出", "hold": "持有"}
        action = action_map.get(t["action"], t["action"])
        stock = t["holding"] if t["holding"] != "空仓" else "-"
        trade_rows.append([
            str(t["date"])[:10], action, stock,
            f"{t['portfolio_value']:,.0f}", ""
        ])
    add_table(slide, headers, trade_rows, 0.65, 1.5, 12, row_h=0.4, col_widths=[2.0, 1.2, 1.8, 2.5, 3.0])

    # ========== Slide 13: 核心问题思考 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "核心问题思考", "为什么理论 Alpha 在真实市场中很难拿到？")
    add_footer(slide, 13)

    lines = [
        "1. 模型假设与现实脱节：A 股受政策、情绪、资金面驱动，短期偏离基本面",
        "",
        "2. PE_TTM 作为日频因子的局限性：",
        "   · PE_TTM 是低频指标（季度更新），日度收益受高频因素驱动",
        "   · 回归 R^2 提升 < 0.3%，PE_TTM 对日度收益几乎无解释力",
        "",
        "3. PEG 策略局限：五年增长率是固定值，无法反映盈利动态变化",
        "",
        "4. 持股集中度风险：全仓单只股票，缺乏分散化",
        "",
        "5. 改进方向：引入动量/波动率因子、动态更新增长率、多股票组合",
    ]
    add_rich_lines(slide, lines, 0.65, 1.5, 12, 5.5, size=13)

    # ========== Slide 14: 结论 ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W, H, fill=COLORS["navy"])
    add_rect(slide, 0, 5.8, W, 0.04, fill=COLORS["blue"])
    add_text(slide, "结论", 0.8, 1.2, 11, 0.7, size=34, bold=True, color=COLORS["paper"], margin=0)

    lines = [
        "1. 二因子模型增量解释力有限：R^2 提升 < 0.3%",
        "",
        "2. PE_TTM 可吸收部分 alpha 异象：白酒三股 alpha 绝对值减少",
        "",
        "3. PEG 策略跑输市场：6.33% vs 27.06%",
        "",
        "4. 低频估值因子难以驱动日频交易 — 核心瓶颈",
    ]
    add_rich_lines(slide, lines, 0.85, 2.2, 11, 4, size=16, color="CBD5E0")

    # Save
    prs.save(PPTX_PATH)
    print(f"PPT saved to {PPTX_PATH.relative_to(ROOT)}")


if __name__ == "__main__":
    build_ppt()
