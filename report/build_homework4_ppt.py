"""Generate a presentation for Homework 4: Multi-Factor Quantitative Stock Selection."""

from __future__ import annotations

import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

import json
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT_DIR = ROOT / "outputs" / "homework4"
PPTX_PATH = OUT_DIR / "homework4_presentation.pptx"
OUT_DIR.mkdir(parents=True, exist_ok=True)

W, H = 13.333, 7.5

COLORS = {
    "ink": "202124",
    "muted": "667085",
    "soft": "F5F7FA",
    "paper": "FFFFFF",
    "green": "1B7F5C",
    "red": "B5403C",
    "navy": "1a1a2e",
    "line": "D8DEE8",
    "blue": "2F6FED",
    "orange": "D97904",
    "purple": "6E56CF",
    "teal": "0D7377",
}


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def add_text(slide, text, x, y, w, h, *, size=18, color=COLORS["ink"],
             bold=False, font="Microsoft YaHei", align=PP_ALIGN.LEFT, margin=0.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_lines(slide, lines, x, y, w, h, *, size=14, color=COLORS["ink"], line_spacing=1.15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(4 * line_spacing)
    return box


def add_rect(slide, x, y, w, h, *, fill=COLORS["paper"], line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_title(slide, title, subtitle=None, *, inverse=False):
    title_color = COLORS["paper"] if inverse else COLORS["ink"]
    sub_color = "CBD5E0" if inverse else COLORS["muted"]
    add_text(slide, title, 0.65, 0.36, 8.6, 0.55, size=26, bold=True, color=title_color, margin=0)
    if subtitle:
        add_text(slide, subtitle, 0.67, 0.95, 9.2, 0.32, size=10, color=sub_color, margin=0)


def add_footer(slide, page, inverse=False):
    c = "CBD5E0" if inverse else "98A2B3"
    add_text(slide, "金融与经济数据挖掘 · 作业4", 0.65, 7.16, 4.2, 0.22, size=8, color=c, margin=0)
    add_text(slide, f"{page:02d}", 12.35, 7.16, 0.35, 0.22, size=8, color=c, align=PP_ALIGN.RIGHT, margin=0)


def add_table(slide, headers, rows, x, y, w, row_h=0.35, col_widths=None, header_color=COLORS["navy"]):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(row_h * n_rows))
    table = table_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw)
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(header_color)
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = "Microsoft YaHei"
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = rgb(COLORS["paper"])
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(COLORS["soft"])
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.size = Pt(9)
                    run.font.color.rgb = rgb(COLORS["ink"])
    return table_shape


def kpi_card(slide, label, value, x, y, color, w=2.8):
    add_rect(slide, x, y, w, 1.4, fill=COLORS["paper"], line=COLORS["line"], radius=True)
    add_text(slide, value, x + 0.15, y + 0.15, w - 0.3, 0.55, size=26, bold=True, color=color, margin=0)
    add_text(slide, label, x + 0.15, y + 0.85, w - 0.3, 0.35, size=11, color=COLORS["muted"], margin=0)


def build_ppt():
    summary_path = OUT_DIR / "summary.json"
    if not summary_path.exists():
        print("summary.json not found, running homework4 main first...")
        from homework4.main import main
        main()
    summary = json.loads(summary_path.read_text(encoding="utf-8"))
    metrics = summary["回测指标"]
    ic_ir = summary["IC_IR"]
    factor_test = summary["单因子检验"]
    weights = summary["因子权重"]

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]

    # ========== Slide 1: Title ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W, H, fill=COLORS["navy"])
    add_rect(slide, 0, 5.8, W, 0.04, fill=COLORS["blue"])
    add_text(slide, "作业4：多因子量化选股实战", 0.8, 1.8, 11, 0.8, size=38, bold=True, color=COLORS["paper"], margin=0)
    add_text(slide, "因子挖掘 → IC/IR质检 → 回归赋权 → 截面选股 → 样本外回测", 0.85, 2.8, 11, 0.5, size=17, color="A0AEC0", margin=0)
    add_text(slide, "上证50成分股  |  SMB + PE_inv + Quality 三因子模型", 0.85, 3.4, 11, 0.4, size=13, color="7B8794", margin=0)
    add_text(slide, "金融与经济数据挖掘", 0.85, 4.5, 6, 0.35, size=13, color="7B8794", margin=0)

    # ========== Slide 2: 实验概述 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "实验概述", "量化投研标准流程完整复现")
    add_footer(slide, 2)

    lines = [
        "一、作业目的",
        "  承接CAPM/二因子模型分析思路，完整复现量化投研标准流程",
        "",
        "二、核心流程（5个模块）",
        "  模块1: 数据获取与预处理 → 模块2: 单因子有效性检验 → 模块3: IC/IR因子质检",
        "  → 模块4: 多因子截面回归赋权 → 模块5: 截面选股+样本外回测",
        "",
        "三、研究设计",
        f"  标的: 上证50成分股（49-50只）",
        f"  训练期: {summary['训练期']}",
        f"  回测期: {summary['回测期']}",
        "  工具: Python + baostock + statsmodels",
    ]
    add_rich_lines(slide, lines, 0.65, 1.5, 12, 5.5, size=14)

    # ========== Slide 3: 因子体系 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "因子体系设计", "三大类因子：市场 / 规模 / 价值 / 质量")
    add_footer(slide, 3)

    headers = ["因子", "符号", "计算方式", "经济含义", "用途"]
    rows = [
        ["市场因子", "MKT", "上证指数月度超额收益率", "系统性风险暴露", "CAPM回归"],
        ["规模因子", "SMB", "总股本 × 月末收盘价", "大盘/小盘风格偏好", "选股因子"],
        ["价值因子", "PE_inv", "1 / PE_TTM（负PE→0）", "低估值→高因子值", "选股因子"],
        ["质量因子", "Quality", "ROE × 股利分红率 × 净利润增速", "高盈利+分红+增长", "选股因子"],
    ]
    add_table(slide, headers, rows, 0.65, 1.5, 12, row_h=0.5, col_widths=[1.2, 1.0, 3.2, 2.5, 1.8])

    lines = [
        "无风险利率 Rf: 月度 0.15%（年化约1.8%）",
        "数据来源: baostock 月度K线 + 年度财务（利润/成长/分红）",
    ]
    add_rich_lines(slide, lines, 0.65, 4.4, 12, 2, size=12, color=COLORS["muted"])

    # ========== Slide 4: 数据预处理 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "模块1：数据获取与预处理", "月度K线 + 财务数据 → 因子计算 → 清洗 → 拆分")
    add_footer(slide, 4)

    lines = [
        "1. 数据获取",
        "   · 月度K线: baostock日频收盘价+PE_TTM → 月末重采样（PE_TTM月度不可用）",
        "   · 财务数据: Q4利润表(ROE/净利润/总股本) + 年增长数据 + 年分红数据",
        "   · 年份匹配: 年度财务数据 → 当年各月（前值填充）",
        "",
        "2. 因子计算",
        "   · SMB = totalShare × close（月末总市值）",
        "   · PE_inv = 1 / PE_TTM（负PE→0，正向化）",
        "   · Quality = ROE × div_ratio × profit_growth",
        "",
        "3. 数据清洗",
        "   · 缺失值: 财务前值填充 + 截面均值填充",
        "   · 异常值: 3σ缩尾（μ±3σ），按月截面独立处理",
        f"   · 训练集: 4,561条/49只股票  |  回测集: 1,131条/50只股票",
    ]
    add_rich_lines(slide, lines, 0.65, 1.5, 12, 5.8, size=12)

    # ========== Slide 5: CAPM筛选 + 单因子检验 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "模块2：单因子有效性检验", "CAPM个股β筛选 → 单因子横截面回归")
    add_footer(slide, 5)

    # CAPM summary
    add_text(slide, "CAPM个股筛选: Ri-Rf = α + β×MKT + ε", 0.65, 1.5, 6, 0.35, size=14, bold=True, color=COLORS["ink"], margin=0)
    add_text(slide, "β显著(p<0.05): 42/49只 (85.7%) → 全部保留进入因子检验", 0.65, 1.85, 8, 0.3, size=12, color=COLORS["muted"], margin=0)

    add_text(slide, "单因子横截面回归: Return(t+1) = α + β×Factor(t) + ε", 0.65, 2.35, 8, 0.35, size=14, bold=True, color=COLORS["ink"], margin=0)

    headers = ["因子", "β均值", "β标准差", "p均值", "t均值", "有效月数", "判定"]
    rows = []
    for f in ["SMB", "PE_inv", "Quality"]:
        info = factor_test[f]
        rows.append([
            {"SMB": "SMB(规模)", "PE_inv": "PE_inv(价值)", "Quality": "Quality(质量)"}[f],
            f"{info['beta_mean']:.6f}",
            f"{info['beta_std']:.6f}",
            f"{info['p_mean']:.4f}",
            f"{info['t_mean']:.4f}",
            str(info['n_months']),
            "✗" if str(info['valid']) == "False" else "✓",
        ])
    add_table(slide, headers, rows, 0.65, 2.8, 12, row_h=0.45, col_widths=[1.8, 1.6, 1.6, 1.3, 1.3, 1.2, 0.9])

    add_text(slide, "结论: 三因子均未通过单因子显著性检验(p>0.05)，按作业要求全部纳入后续分析", 0.65, 4.8, 10, 0.3, size=11, color=COLORS["red"], margin=0)

    # ========== Slide 6: IC/IR因子质检 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "模块3：IC/IR因子质检", "Z-Score标准化 → 月度截面Pearson IC → IR计算")
    add_footer(slide, 6)

    headers = ["因子", "IC均值", "IC标准差", "IR", "评级", "IR判定", "IC阈值"]
    rows = []
    for f in ["SMB", "PE_inv", "Quality"]:
        info = ic_ir[f]
        rows.append([
            {"SMB": "SMB(规模)", "PE_inv": "PE_inv(价值)", "Quality": "Quality(质量)"}[f],
            f"{info['IC_mean']:.4f}",
            f"{info['IC_std']:.4f}",
            f"{info['IR']:.4f}",
            info['grade'],
            "✓" if abs(info['IR']) > 0.1 else "✗",
            ">0.05优秀, >0.02预测" if f == "Quality" else "",
        ])
    add_table(slide, headers, rows, 0.65, 1.6, 12, row_h=0.5, col_widths=[1.8, 1.3, 1.3, 1.3, 1.5, 1.0, 2.5])

    # Insert IC/IR chart
    chart_path = OUT_DIR / "ic_ir_summary.png"
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), Inches(0.4), Inches(3.4), Inches(10.5), Inches(3.5))

    # ========== Slide 7: IC序列 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "因子IC序列（月度截面）", "训练期2015-2023各月IC动态")
    add_footer(slide, 7)

    ic_path = OUT_DIR / "ic_series.png"
    if ic_path.exists():
        slide.shapes.add_picture(str(ic_path), Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.2))

    lines = [
        "Quality因子IC均值最高(0.023)，具备预测能力 | SMB因子IC为负(-0.034)暗示大盘股溢价",
        "三因子IC标准差均较大(0.19-0.26)，月度IC波动剧烈，因子预测稳定性有限",
    ]
    add_rich_lines(slide, lines, 0.65, 6.7, 12, 0.6, size=11, color=COLORS["muted"])

    # ========== Slide 8: 多因子回归赋权 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "模块4：多因子回归静态赋权", "Return(t+1) = α + w1×SMB_std + w2×PE_inv_std + w3×Quality_std + ε")
    add_footer(slide, 8)

    headers = ["参数", "系数", "p值", "显著性", "N", "R²"]
    rows = [
        ["截距 α", f"{summary['回归截距']:.6f}", "0.0000", "***", "", ""],
        ["w1 (SMB)", f"{weights['SMB']:.6f}", "0.0033", "**", "", ""],
        ["w2 (PE_inv)", f"{weights['PE_inv']:.6f}", "0.9007", "不显著", "", ""],
        ["w3 (Quality)", f"{weights['Quality']:.6f}", "0.2860", "不显著", "4,561", f"{summary['回归R2']:.4f}"],
    ]
    add_table(slide, headers, rows, 0.65, 1.5, 10, row_h=0.45, col_widths=[1.8, 1.6, 1.2, 1.2, 1.2, 1.2])

    lines = [
        "综合得分: Score = -0.00547×SMB_std + 0.00023×PE_inv_std + 0.00185×Quality_std",
        "· SMB权重为负 → 小盘股得分更高  ·  Quality权重为正 → 高ROE+高分红+高增长得分更高",
        "· PE_inv权重≈0 → 估值因子在控制规模和质量的条件下几乎无边际贡献",
        "· 模型R²=0.24% → 因子对收益线性解释力低（月度横截面回归常见现象）",
    ]
    add_rich_lines(slide, lines, 0.65, 3.8, 12, 3.5, size=12)

    # Factor weights chart
    weights_path = OUT_DIR / "factor_weights.png"
    if weights_path.exists():
        slide.shapes.add_picture(str(weights_path), Inches(8.5), Inches(1.5), Inches(4.2), Inches(2.8))

    # ========== Slide 9: 截面选股策略 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "模块5：截面打分选股", "月度Top3选股 + 等权配置 + 月度调仓")
    add_footer(slide, 9)

    lines = [
        "选股规则:",
        "  · 每月月末按综合得分排序，选取Top 3",
        "  · 无法交易时顺延至下一名",
        "",
        "交易规则:",
        "  · 调仓频率: 每月1次（月末调仓）",
        "  · 仓位管理: 3只等权重，单票上限 ≤ 40%",
        "  · 交易成本: 单边0.3%（买卖双向）",
        "  · 初始资金: 100万元",
        "",
        "关键原则: 无前视偏差（Look-ahead Bias Free）",
        "  月末因子值 → 预测下月收益 → 月末调仓 → 下月收益归属",
    ]
    add_rich_lines(slide, lines, 0.65, 1.5, 6.5, 5.5, size=14)

    # Scoring example
    add_text(slide, "首次选股示例 (2024-01)", 7.5, 1.5, 5, 0.3, size=13, bold=True, color=COLORS["ink"], margin=0)
    headers = ["排名", "代码", "得分"]
    rows = [
        ["1", "sh.688008 澜起科技", "0.01166"],
        ["2", "sh.603993 洛阳钼业", "0.01162"],
        ["3", "sh.603501 豪威集团", "0.00838"],
        ["4", "sh.600111 北方稀土", "0.00794"],
        ["5", "sh.688041 海光信息", "0.00728"],
    ]
    add_table(slide, headers, rows, 7.5, 2.0, 5.2, row_h=0.4, col_widths=[0.6, 2.8, 1.5])

    # ========== Slide 10: 回测结果 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "回测结果总览", "2024-2025 样本外回测 | 策略 vs 上证指数")
    add_footer(slide, 10)

    kpi_card(slide, "策略累计收益", f"{metrics['累计收益率']:.2%}", 0.65, 1.5, COLORS["blue"])
    kpi_card(slide, "上证指数累计收益", f"{metrics['上证指数累计收益']:.2%}", 3.7, 1.5, COLORS["orange"])
    kpi_card(slide, "超额收益", f"{metrics['超额收益(vs上证指数)']:.2%}", 6.75, 1.5, COLORS["green"])
    kpi_card(slide, "策略最大回撤", f"{metrics['最大回撤']:.2%}", 9.8, 1.5, COLORS["red"])

    kpi_card(slide, "年化收益率", f"{metrics['年化收益率']:.2%}", 0.65, 3.15, COLORS["teal"])
    kpi_card(slide, "月度胜率", f"{metrics['月度胜率']:.1%}", 3.7, 3.15, COLORS["purple"])
    kpi_card(slide, "夏普比率", f"{metrics['夏普比率']:.4f}", 6.75, 3.15, COLORS["navy"])
    kpi_card(slide, "卡玛比率", f"{metrics['卡玛比率']:.4f}", 9.8, 3.15, COLORS["ink"])

    # NAV chart
    nav_path = OUT_DIR / "nav_vs_market.png"
    if nav_path.exists():
        slide.shapes.add_picture(str(nav_path), Inches(0.4), Inches(4.8), Inches(12.5), Inches(2.4))

    # ========== Slide 11: 回撤分析 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "回撤分析", "策略回撤 vs 市场回撤")
    add_footer(slide, 11)

    dd_path = OUT_DIR / "drawdown.png"
    if dd_path.exists():
        slide.shapes.add_picture(str(dd_path), Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.5))

    lines = [
        f"策略最大回撤: {metrics['最大回撤']:.2%}  |  上证指数最大回撤: {metrics['上证指数最大回撤']:.2%}",
        "策略回撤主要出现在2024年8月和2025年10-11月，与市场回调同步但幅度更大",
        "2024年9月(+20.5%)和2025年8月(+26.9%)的大幅反弹快速修复了前期回撤",
    ]
    add_rich_lines(slide, lines, 0.65, 6.2, 12, 1.2, size=12, color=COLORS["muted"])

    # ========== Slide 12: 持仓与调仓分析 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "持仓分析与调仓记录", "核心持仓集中度分析")
    add_footer(slide, 12)

    headers = ["股票代码", "股票名称", "被选次数", "占比", "月均收益"]
    rows = [
        ["sh.688008", "澜起科技", "22/22", "100%", "+4.13%"],
        ["sh.603993", "洛阳钼业", "11/22", "50%", "+6.17%"],
        ["sh.603501", "豪威集团", "11/22", "50%", "+1.52%"],
        ["sh.601127", "赛力斯", "8/22", "36%", "+3.58%"],
        ["sh.600111", "北方稀土", "5/22", "23%", "+4.87%"],
        ["sh.603259", "药明康德", "4/22", "18%", "+1.52%"],
        ["sh.600031", "三一重工", "3/22", "14%", "+1.67%"],
    ]
    add_table(slide, headers, rows, 0.65, 1.5, 7, row_h=0.45, col_widths=[1.4, 1.4, 1.2, 1.0, 1.2])

    add_text(slide, "月度调仓记录（前7月）", 8.2, 1.5, 5, 0.3, size=13, bold=True, color=COLORS["ink"], margin=0)
    trade_log = pd.read_csv(OUT_DIR / "trade_log.csv")
    headers2 = ["月份", "持仓", "月收益"]
    trade_rows = []
    for _, t in trade_log.head(7).iterrows():
        trade_rows.append([
            str(t["month"])[:7],
            t["holdings"][:35] + ("..." if len(str(t["holdings"])) > 35 else ""),
            f"{t['month_return']:.2%}",
        ])
    add_table(slide, headers2, trade_rows, 8.2, 2.0, 4.8, row_h=0.45, col_widths=[1.0, 2.8, 0.8])

    add_text(slide, "澜起科技(内存接口芯片)全程持有，是策略核心收益来源；AI/半导体产业链景气上行驱动", 0.65, 5.5, 12, 0.3, size=11, color=COLORS["muted"], margin=0)

    # ========== Slide 13: 年度表现 ==========
    slide = prs.slides.add_slide(blank)
    add_title(slide, "分年度表现", "2024 vs 2025 收益分解")
    add_footer(slide, 13)

    headers = ["年份", "策略收益", "上证指数收益", "超额收益", "策略最大回撤", "正收益月数"]
    rows = [
        ["2024", "+27.31%", "+20.20%", "+7.11%", "-15.48%", "7/11 (63.6%)"],
        ["2025", "+87.38%", "+19.63%", "+67.75%", "-13.84%", "7/11 (63.6%)"],
    ]
    add_table(slide, headers, rows, 0.65, 1.5, 12, row_h=0.6, col_widths=[1.2, 2.0, 2.0, 2.0, 2.0, 2.0])

    lines = [
        "2024年:",
        "  上半年震荡蓄力，净值在97万-130万间波动",
        "  9月(+20.5%)大反弹突破，10-12月小幅回调",
        "  全年+27.31%，超越上证+7.11%",
        "",
        "2025年:",
        "  7-9月主升浪，三个月累计涨幅超70%",
        "  8月单月+26.9%为策略最佳月份",
        "  10-11月回撤-13.8%，年末净值收于244万",
        "  全年+87.38%，大幅跑赢上证+67.75%",
    ]
    add_rich_lines(slide, lines, 0.65, 3.2, 12, 4, size=13)

    # ========== Slide 14: 结论 ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W, H, fill=COLORS["navy"])
    add_rect(slide, 0, 5.8, W, 0.04, fill=COLORS["blue"])
    add_text(slide, "结论与展望", 0.8, 1.2, 11, 0.7, size=34, bold=True, color=COLORS["paper"], margin=0)

    lines = [
        "1. Quality因子表现最优: IC=0.023, IR=0.116, 是唯一具备预测能力的因子",
        "",
        "2. 三因子组合样本外效果显著: 策略+144% vs 上证+39%, 超额+105%",
        "",
        "3. 持仓高度集中于澜起科技: 全程22/22月选中, 半导体/AI产业链驱动",
        "",
        "4. 模型R²极低(0.24%): 线性模型的解释力有限, 可引入ML方法提升",
        "",
        "5. 改进方向: 引入动量/波动率因子 → 行业中性化 → 滚动动态权重 → XGBoost/LGBM",
    ]
    add_rich_lines(slide, lines, 0.85, 2.2, 11, 4.5, size=16, color="CBD5E0")

    # ========== Slide 15: Thank you ==========
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W, H, fill=COLORS["navy"])
    add_text(slide, "谢谢！", 0.8, 3.0, 11, 1, size=44, bold=True, color=COLORS["paper"], align=PP_ALIGN.CENTER, margin=0)
    add_text(slide, "金融与经济数据挖掘 · 作业4 · 多因子量化选股实战", 0.8, 4.2, 11, 0.5, size=16, color="A0AEC0", align=PP_ALIGN.CENTER, margin=0)

    # Save
    prs.save(PPTX_PATH)
    print(f"PPT saved to {PPTX_PATH.relative_to(ROOT)}")


if __name__ == "__main__":
    build_ppt()
